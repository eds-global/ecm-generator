import os, sys, time, webbrowser
import pandas as pd
import streamlit as st
import plotly.express as px
from PIL import Image as PILImage
from PIL import Image
import numpy as np
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
import json
import streamlit.components.v1 as components
import pdfplumber
import re
import io
import matplotlib.image as mpimg
from streamlit.runtime.uploaded_file_manager import UploadedFile
from io import BytesIO
import time
import random
from datetime import datetime
from shutil import copyfile
import glob as gb
import subprocess
import shutil
from pathlib import Path
from collections import defaultdict
import streamlit.web.cli as stcli
from src import insertWall, insertConst, orient, lighting, equip, windows, insertRoof, wwr
import traceback
from helper import *
from report_ext import *
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import statsmodels.api as sm 
from src import ModifyWallRoof
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
import plotly.io as pio
from src import lv_b, ls_c, lv_d, pv_a_loop, sv_a, beps, bepu, lvd_summary, sva_zone, locationInfo, masterFile, sva_sys_type, pv_a_pump, pv_a_heater, pv_a_equip, pv_a_tower, ps_e, inp_shgc

# --- Streamlit Page Config ---
st.set_page_config(page_title="AWECM Sim", page_icon="⚡", layout='wide')

# --- Inject Custom CSS ---
st.markdown("""
    <style>
        .block-container { padding-top: 0rem !important; }
        header, main { margin-top: 0rem !important; padding-top: 0rem !important; }
        .stButton>button { box-shadow: 1px 1px 1px rgba(0, 0, 0, 0.8); }
        .heading-with-shadow {
            text-align: left;
            color: red;
            text-shadow: 0px 8px 4px rgba(255, 255, 255, 0.4);
            background-color: white;
        }
        body {
            background-color: #bfe1ff;
            animation: changeColor 5s infinite;
        }
        #MainMenu, footer, .viewerBadge_container__1QSob {visibility: hidden;}
        header .stApp [title="View source on GitHub"] { display: none; }
        .stApp header, .stApp footer {visibility: hidden;}
        .stButton button {
            height: 30px;
            width: 166px;
        }
    </style>
""", unsafe_allow_html=True)

def resource_path(relative_path):
    """Get absolute path to resource (works for dev and PyInstaller exe)"""
    if getattr(sys, 'frozen', False):  # Running in exe
        base_path = sys._MEIPASS
    else:
        base_path = os.path.dirname(__file__)
    return os.path.join(base_path, relative_path)

def remove_utility(inp_data):
    start_marker = "Utility Rates"
    end_marker = "Output Reporting"

    start_index, end_index = None, None

    # Find first occurrence of start_marker
    for i, line in enumerate(inp_data):
        if start_marker in line:
            start_index = i + 3

    # Find first occurrence of end_marker after start_marker
    if start_index is not None:
        for i, line in enumerate(inp_data[start_index:], start=start_index):
            if end_marker in line:
                end_index = i - 3
                break

    if start_index is None or end_index is None:
        raise ValueError("Could not find section markers in INP file.")

    # Remove everything between start_index and end_index (inclusive)
    new_data = inp_data[:start_index] + inp_data[end_index + 1:]
    return new_data

def remove_betweenLightEquip(inp_data):
    start_marker = "Floors / Spaces / Walls / Windows / Doors"
    end_marker = "Electric & Fuel Meters"

    start_index, end_index = None, None

    # Find first occurrence of start_marker
    for i, line in enumerate(inp_data):
        if start_marker in line:
            start_index = i + 2
            break

    # Find first occurrence of end_marker after start_marker
    if start_index is not None:
        for i, line in enumerate(inp_data[start_index:], start=start_index):
            if end_marker in line:
                end_index = i - 2
                break

    if start_index is None or end_index is None:
        raise ValueError("Could not find section markers in INP file.")

    new_data = []
    cleaning = False  # to track inside LIGHTING/EQUIPMENT block

    for i, line in enumerate(inp_data):
        if start_index <= i <= end_index:
            if line.strip().startswith("LIGHTING-W/AREA") or line.strip().startswith("EQUIPMENT-W/AREA"):
                cleaning = True
                new_data.append(line)  # keep the main line
                continue
            if cleaning:
                if line.strip().startswith("*"):
                    continue  # skip continuation lines
                else:
                    cleaning = False  # stop cleaning when normal line comes
        new_data.append(line)
    return new_data

if "script_choice" not in st.session_state:
    st.session_state.script_choice = "home"

if "tools_dropdown" not in st.session_state:
    st.session_state.tools_dropdown = "Select"

if "reset_tools" not in st.session_state:
    st.session_state.reset_tools = False

# --- Header Section ---

# 🔁 Reset dropdown BEFORE widget is created
if st.session_state.reset_tools:
    st.session_state.tools_dropdown = "Select"
    st.session_state.reset_tools = False

col1, col2, col3, _, col4, col5 = st.columns([0.16,0.16,2.3,0.4,0.5,0.5])

with col1:
    st.image("images/analysis.png")

with col2:
    st.image("images/eds.png")

with col3:
    st.markdown("""
    <div style="text-align:center;">
        <h3 style="margin-bottom:2px; color:black;">
            <span style="color: rgb(202, 50, 50);">
                Automating Workflows for Energy Simulation
            </span>
        </h3>
    </div>
    """, unsafe_allow_html=True)

with col4:
    tool = st.selectbox(
        "",
        ["ParSim", "ComSim"],
        key="tools_dropdown",
        label_visibility="collapsed"
    )

    if tool == "ParSim":
        st.session_state.script_choice = "tool1"

    elif tool == "ComSim":
        st.session_state.script_choice = "tool2"

with col5:
    if st.button("Home", key="home"):
        st.session_state.script_choice = "home"
        st.session_state.reset_tools = True   # ✅ trigger reset

# Reduce gap before red line
st.markdown("""<div style='margin-top:-40px;'><hr style="border:1px solid red"></div>""",unsafe_allow_html=True)

# ---------------------- HEADER STYLES -------------------------
st.markdown("""<style></style>""", unsafe_allow_html=True)
# # ---------------------- HEADER -------------------------
# st.markdown("""
# <div class="header-container">
#     <div class="header-title">Automating Workflows for Energy Simulation</div>
#     <div class="header-buttons">
#         <a href="#"><button class="css-1emrehy edgvbvh10">Tools</button></a>
#         <a href="#"><button class="css-1emrehy edgvbvh10">Home</button></a>
#     </div>
# </div>
# """, unsafe_allow_html=True)

if st.session_state.script_choice == "home":
    # ---------------------- THREE COLUMN CARDS -------------------------
    col1, col2, col3 = st.columns(3)
    # --------- Card 1: AWESim ----------
    with col1:
        st.image("images/awesim.png")
    # --------- Card 2: PARSim ----------
    with col2:
        col1, col2 = st.columns(2)
        with col1:
            st.image("images/6.png")
        with col2:
            if st.button("ParSim"):
                st.session_state.script_choice = "tool1"
                
    # --------- Card 3: COMSim ----------
    with col3:
        col1, col2 = st.columns(2)
        with col1:
            st.image("images/8.png")
        with col2:
            if st.button("ComSim", key="1"):
                st.session_state.script_choice = "tool2"

    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown(
            '<span style="font-size:24px; font-weight:600;">About <span style="color: rgb(202,50,50);">AWESim</span></span>',
            unsafe_allow_html=True
        )
        # st.markdown('<span style="font-size:24px; font-weight:600;">About <span style="color:red;">AWESim</span></span>',unsafe_allow_html=True)
        st.write("""AWESim enables an energy simulation expert to discover optimal design parameters through systematic parametric investigations in just a few clicks. Please share your feedback with us.
        These investigations are accessible through interactive charts and downloadable reports. Please feel free to share your feedback with us at **support@edsglobal.com**
        """)
    with col2:
        st.markdown('<span style="font-size:24px; font-weight:600;">Parametric Analysis</span>',unsafe_allow_html=True)
        st.write("""AWESim enables systematic parametric investigations with ease. Interactive charts and downloadable reports provide clear insights. 
        These investigations are accessible through interactive charts and downloadable reports. We keep improving this tool — please send feedback.""")
    with col3:
        st.markdown('<span style="font-size:24px; font-weight:600;">ECM Comparison</span>',unsafe_allow_html=True)
        st.write("""AWESim enables quick ECM comparison through interactive charts and downloadable reports. Please share your feedback with us. These investigations are accessible through interactive charts and downloadable reports. Please share your feedback with us.""")

    # ---------------------- THREE COLUMN CARDS -------------------------
    col1, col2, col3 = st.columns(3)
    # --------- Card 1: AWESim ----------
    with col1:
        st.markdown('<span style="font-size:24px; font-weight:600;">About <span style="color: rgb(202,50,50);">EDS</span></span>',unsafe_allow_html=True)
        st.write("""Environmental Design Solutions [EDS] is a sustainability advisory firm. Since 2002, EDS has worked on over 500 green building and energy efficiency projects worldwide.
        The team focuses on climate change mitigation, low-carbon design, building simulation, performance audits, and capacity building.""")
        st.markdown("</div>", unsafe_allow_html=True)
    # --------- Card 2: PARSim ----------
    with col2:
        st.markdown('<span style="font-size:24px; font-weight:600;">This Version</span>',unsafe_allow_html=True)
        st.write("""
        **Build v1.0.0**

        **Fixes**-
        No Fixes right now!  
        """)
        st.markdown("</div>", unsafe_allow_html=True)
    # --------- Card 3: COMSim ----------
    with col3:
        st.markdown('<span style="font-size:24px; font-weight:600;">Disclaimer & Acknowledgement</span>',unsafe_allow_html=True)
        st.write("""AWEPCM Sim is the outcome of best efforts by simulation experts at EDS. EDS does not assume responsibility for outcomes from this application. ECM comparison through interactive charts and downloadable reports. The user indemnifies EDS of any damages.  
        (more)""")
        st.markdown("</div>", unsafe_allow_html=True)

if st.session_state.script_choice == "tool1":
    # Load location database
    csv_path = resource_path(os.path.join("database", "Simulation_locations.csv"))
    weather_df = pd.read_csv(csv_path)
    db_path = resource_path(os.path.join("database", "AllData.xlsx"))
    output_csv = resource_path(os.path.join("database", "Randomized_Sheet.xlsx"))
    location_path = resource_path(os.path.join("database", "Simulation_locations.csv"))
    location = pd.read_csv(location_path)
    locations = sorted(location['Sim_location'].unique())
    updated_df = pd.read_excel(output_csv)

    updated_df["Wall_Roof_Glazing_WWR_Orient_Light_Equip"] = (
        updated_df["Wall"].astype(str) + "_" +
        updated_df["Roof"].astype(str) + "_" +
        updated_df["Glazing"].astype(str) + "_" +
        updated_df["WWR"].astype(str) + "_" +
        updated_df["Orient"].astype(str) + "_" +
        updated_df["Light"].astype(str) + "_" +
        updated_df["Equip"].astype(str)
    )
    # Inputs
    col1, col2, col3, col4 = st.columns(4)
    st.markdown("""
        <style>
        div[data-testid="stFileUploader"] section {
            padding: 0.01rem 0 !important;  /* Reduce vertical padding */
        }
        div[data-testid="stFileUploader"] div[role="button"] {
            padding: 0.0rem 0.0rem !important;  /* Reduce height of clickable area */
            font-size: 0.00rem !important;      /* Smaller text */
        }
        </style>
    """, unsafe_allow_html=True)
    # --- NEW: handle both Country and Location ---
    if "Country" in location.columns:
        # If CSV already has Country column
        countries = sorted(location["Country"].unique().tolist())
    else:
        # If your current CSV has only Indian cities — default to India
        countries = ["India"]
    with col1:
        st.write("📝 Project Name")
        project_name = st.text_input("", placeholder="Enter project name", label_visibility="collapsed")
        # project_name = st.text_input("📝 Project Name", placeholder="Enter project name")
        project_name_clean = project_name.replace(" ", "_")
        user_nm = project_name_clean
        if project_name_clean:
            parent_dir = os.path.dirname(os.getcwd())
            batch_outputs_dir = os.path.join(parent_dir, "Batch_Outputs")
            project_folder = os.path.join(batch_outputs_dir, project_name_clean)
            # Check if project folder already exists
            if os.path.exists(project_folder):
                st.warning("⚠️ Project name already exists! Please select another name.")
            #     st.stop()
    with col2:
        # Add "Other" option to countries
        countries.append("Custom Weather")
        st.write("🌎 Select Country")
        selected_country = st.selectbox("", countries, label_visibility="collapsed")

        # Filter and sort locations for the selected country (if not "Other")
        if selected_country != "Custom Weather":
            if "Country" in location.columns:
                filtered_locations = (
                    location[location["Country"] == selected_country]["Sim_location"]
                    .dropna()
                    .unique()
                    .tolist()
                )
                filtered_locations = sorted(filtered_locations)
            else:
                filtered_locations = sorted(location["Sim_location"].dropna().tolist())
        else:
            filtered_locations = []  # No locations for "Other"

    bin_name = ""
    # Only show City dropdown if not "Other"
    if selected_country != "Custom Weather":
        with col3:
            st.write("🌎 Select City")
            user_input = st.selectbox("", filtered_locations, label_visibility="collapsed").lower()
    else:
        with col3:
            user_input = "Other-City"
            # When "Other" is selected, show .bin upload option
            st.write("📤 Upload .bin file")
            uploaded_bin = st.file_uploader("", type=["bin"], label_visibility="collapsed")
            if uploaded_bin is not None:
                save_folder = r"C:\doe22\weather"
                os.makedirs(save_folder, exist_ok=True)
                # Always rename uploaded file to 1.bin
                save_path = os.path.join(save_folder, "1.bin")
                # Save uploaded file as 1.bin
                with open(save_path, "wb") as f:
                    f.write(uploaded_bin.getbuffer())
                bin_name = "1"   # without extension

            # st.markdown(f"<span style='color:red;'>Remember that uploaded file in your C:/doe22/weather</span>", unsafe_allow_html=True)
            # if uploaded_bin is not None:
            #     # Define the save folder path
            #     save_folder = r"C:\doe22\weather"
            #     os.makedirs(save_folder, exist_ok=True)

            #     # Original filename
            #     original_name = uploaded_bin.name
            #     name_without_ext, ext = os.path.splitext(original_name)

            #     # Default save path
            #     save_path = os.path.join(save_folder, original_name)

            #     # Check if file already exists → create copy versions
            #     counter = 1
            #     while os.path.exists(save_path):
            #         new_name = f"{name_without_ext}_copy{counter}{ext}"
            #         save_path = os.path.join(save_folder, new_name)
            #         counter += 1

            #     # Save the file
            #     with open(save_path, "wb") as f:
            #         f.write(uploaded_bin.getbuffer())
            #     bin_name = name_without_ext
                # st.success(f"File saved as: {os.path.basename(save_path)}")
            
    with col4:
        uploaded_file = st.file_uploader("📤 Upload eQUEST INP file", type=["inp"])
        if uploaded_file:
            # st.write(uploaded_file)
            if uploaded_file.name != 'F1.inp':
                uploaded_file.name = user_nm + '.inp'

            # Go one step outside current working directory
            parent_dir = os.path.dirname(os.getcwd())

            # Create path for Batch_Outputs folder
            batch_outputs_dir = os.path.join(parent_dir, "Batch_Outputs")

            # Make sure the folder exists
            os.makedirs(batch_outputs_dir, exist_ok=True)

            # Save uploaded file in Batch_Outputs
            uploaded_inp_path = os.path.join(batch_outputs_dir, uploaded_file.name)
            with open(uploaded_inp_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            # Output folder = same folder where uploaded file is saved
            output_inp_folder = os.path.dirname(uploaded_inp_path)
            inp_folder = output_inp_folder

            project_folder = os.path.join(batch_outputs_dir, project_name_clean)

    run_cnt = 1
    location_id, weather_path = "", ""
    if user_input != "Other-City":
        matched_row = weather_df[weather_df['Sim_location'].str.lower().str.contains(user_input)]
        # st.write(matched_row)
    else:
        matched_row = pd.DataFrame()

    if not matched_row.empty:
        location_id = matched_row.iloc[0]['Location_ID']
        weather_path = matched_row.iloc[0]['Weather_file']
    elif user_input:
        weather_path = bin_name
        # st.error("❌ Location not found.")
    # st.write(weather_path)
    # Start simulation
    if st.button("Simulate 🚀"):
        # st.clear_cache()
        if uploaded_file is None:
            st.warning("⚠️ Please Upload .INP File!")
            st.stop()
        if bin_name is None:
            st.warning("⚠️ Please Upload .BIN File!")
            st.stop()
        if not project_name_clean:
            st.warning("⚠️ Please enter a project name.")
            st.stop()
        with st.spinner("⚡ Processing... This may take a few minutes."):
            os.makedirs(output_inp_folder, exist_ok=True)
            new_batch_id = f"{int(time.time())}"  # unique ID

            selected_rows = updated_df[updated_df['Batch_ID'] == run_cnt]
            batch_output_folder = os.path.join(output_inp_folder, f"{user_nm}")
            os.makedirs(batch_output_folder, exist_ok=True)

            num = 1
            modified_files = []
            for _, row in selected_rows.iterrows():
                selected_inp = uploaded_file.name
                new_inp_name = f"{row['Wall']}_{row['Roof']}_{row['Glazing']}_{row['Orient']}_{row['Light']}_{row['WWR']}_{row['Equip']}_{selected_inp}"
                new_inp_path = os.path.join(batch_output_folder, new_inp_name)

                inp_file_path = os.path.join(inp_folder, selected_inp)
                if not os.path.exists(inp_file_path):
                    st.error(f"File {inp_file_path} not found. Skipping.")
                    continue

                # st.info(f"Modifying INP file {num}: {selected_inp} -> {new_inp_name}")
                num += 1

                # Apply modifications
                inp_content = wwr.process_window_insertion_workflow(inp_file_path, row["WWR"])
                # inp_content = orient.updateOrientation(inp_content, row["Orient"])
                inp_content = lighting.updateLPD(inp_content, row['Light'])
                # inp_content = insertWall.update_Material_Layers_Construction(inp_content, row["Wall"])
                # inp_content = insertRoof.update_Material_Layers_Construction(inp_content, row["Roof"])
                # inp_content = insertRoof.removeDuplicates(inp_content)
                inp_content = equip.updateEquipment(inp_content, row['Equip'])
                inp_content = windows.insert_glass_types_multiple_outputs(inp_content, row['Glazing'])
                inp_content =remove_utility(inp_content)
                if row['Light'] > 0:
                    inp_content =remove_betweenLightEquip(inp_content)
                count = ModifyWallRoof.count_exterior_walls(inp_content)
                if count > 1:
                    inp_content = ModifyWallRoof.fix_walls(inp_content, row["Wall"])
                    inp_content = ModifyWallRoof.fix_roofs(inp_content, row["Roof"])
                    # inp_content = insertRoof.removeDuplicates(inp_content)
                    with open(new_inp_path, 'w') as file:
                        file.writelines(inp_content)
                    modified_files.append(new_inp_name)
                else:
                    st.write("No Exterior-Wall Exists!")
        
            simulate_files = []
            # Copy and run batch script
            if uploaded_file is None:
                st.error("Please upload an INP file before starting the simulation.")
            else:
                st.markdown(f"<span style='color:green;'>✅ Updating DAYLIGHTING from YES to NO!</span>", unsafe_allow_html=True)
                script_dir = os.path.dirname(os.path.abspath(__file__))
                shutil.copy(os.path.join(script_dir, "script.bat"), batch_output_folder)
                inp_files = [f for f in os.listdir(batch_output_folder) if f.lower().endswith(".inp")]
                for inp_file in inp_files:
                    file_path = os.path.join(batch_output_folder, os.path.splitext(inp_file)[0])
                    subprocess.call(
                        [os.path.join(batch_output_folder, "script.bat"), file_path, weather_path],
                        shell=True
                    )
                    simulate_files.append(inp_file)
            
                subprocess.call([os.path.join(batch_output_folder, "script.bat"), batch_output_folder, weather_path], shell=True)
                # if os.path.exists(save_path):
                #     os.remove(save_path)
                required_sections = ['BEPS', 'BEPU', 'LS-C', 'LV-B', 'LV-D', 'PS-E', 'SV-A']
                log_file_path = check_missing_sections(batch_output_folder, required_sections, new_batch_id, user_nm)
                get_failed_simulation_data(batch_output_folder, log_file_path)
                clean_folder(batch_output_folder)
                combined_Data = get_files_for_data_extraction(batch_output_folder, log_file_path, new_batch_id, location_id, user_nm, user_input)
            logFile = pd.read_excel(log_file_path)
            total_runs = len(updated_df)
            total_sims = len(logFile)
            success_count = (logFile["Status"] == "Success").sum()
            success_rate = (success_count / total_sims) * 100 if total_sims > 0 else 0
        # exportCSV = resource_path(os.path.join("Export.csv"))
        # combined_Data = pd.read_csv(exportCSV)
        combined_Data["Equip(W/Sqft)"] = combined_Data["Equipment-Total(W)"] / combined_Data["Floor-Total-Above-Grade(SQFT)"]
        combined_Data["Light(W/Sqft)"] = combined_Data["Power Lighting Total(W)"] / combined_Data["Floor-Total-Above-Grade(SQFT)"]
        wall_to_rvalue = {
            1: 2.5, 2: 5.0, 3: 7.5, 4: 10.0,
            5: 12.5, 6: 15.0, 7: 17.5, 8: 20.0,
            9: 22.5, 10: 25.0, 11: 27.5, 12: 30.0
        }

        # Extract wall and roof codes
        combined_Data["WallCode"] = combined_Data["FileName"].str.split("_").str[0].astype(int)
        combined_Data["RoofCode"] = combined_Data["FileName"].str.split("_").str[1].astype(int)

        # Map to R-values
        combined_Data["R-Value-Wall"] = combined_Data["WallCode"].map(wall_to_rvalue)
        combined_Data["R-Value-Roof"] = combined_Data["RoofCode"].map(wall_to_rvalue)
        # st.write(combined_Data)
        # Drop intermediate columns if you don’t need them
        combined_Data.drop(columns=["WallCode", "RoofCode"], inplace=True)

        # Split FileName into parts
        split_cols = combined_Data['FileName'].str.split("_", expand=True)

        # Assign names to first 8 parts
        split_cols.columns = ["wall", "roof", "glazing", "orient", "light", "wwr", "equip", "suffix"]

        # Convert numeric columns (except suffix)
        for col in split_cols.columns[:-1]:
            split_cols[col] = pd.to_numeric(split_cols[col], errors="coerce")

        # Merge back
        combined_Data_expanded = pd.concat([combined_Data, split_cols], axis=1)

        # Filter subsets
        wall_df    = combined_Data_expanded[(combined_Data_expanded["wall"] > 0) | (combined_Data_expanded.index == 0)]
        roof_df    = combined_Data_expanded[(combined_Data_expanded["roof"] > 0) | (combined_Data_expanded.index == 0)]
        glazing_df = combined_Data_expanded[(combined_Data_expanded["glazing"] > 0) | (combined_Data_expanded.index == 0)]
        wwr_df     = combined_Data_expanded[(combined_Data_expanded["wwr"] > 0) | (combined_Data_expanded.index == 0)]
        orient_df  = combined_Data_expanded[(combined_Data_expanded["orient"] > 0) | (combined_Data_expanded.index == 0)]
        light_df   = combined_Data_expanded[(combined_Data_expanded["light"] > 0) | (combined_Data_expanded.index == 0)]
        equip_df   = combined_Data_expanded[(combined_Data_expanded["equip"] > 0) | (combined_Data_expanded.index == 0)]

        wall_roof_df = pd.concat([
            wall_df.assign(Variable="Wall", X=wall_df["N-Wall-U-Value(BTU/HR-SQFT-F)"]),
            roof_df.assign(Variable="Roof", X=roof_df["ROOF-U-Value(BTU/HR-SQFT-F)"])
        ])

        light_equip_df = pd.concat([
            light_df.assign(Variable="Lighting", X=light_df["Light(W/Sqft)"]),
            equip_df.assign(Variable="Equipment", X=equip_df["Equip(W/Sqft)"])
        ])

        glazing_df = pd.concat([
            glazing_df.assign(Variable="Glazing", X=glazing_df["SC"])
        ])

        # Common y-axis range
        y_min = combined_Data_expanded["Energy_Outcome(KWH)"].min()
        y_max = combined_Data_expanded["Energy_Outcome(KWH)"].max()

        def make_savings_wwr_barplot(df, title="Energy Savings vs WWR", x_label="Window-to-Wall Ratio"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to Baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Melt EUI & WWR into one column for plotting (if multiple types, else just use WWR)
            df_melt = pd.melt(
                df,
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["WWR"],  # replace/add if multiple types
                var_name="Variable",
                value_name="WWR_Value"
            )

            # Plot
            fig = px.bar(
                df_melt.iloc[1:],  # exclude baseline
                x="WWR_Value",
                y="PercentSaving",
                color="Variable",
                barmode="group",
                text="PercentSaving",
                labels={
                    "WWR_Value": x_label,
                    "PercentSaving": "% Saving",
                    "Variable": "Component"
                },
                color_discrete_map={"WWR": "blue"}  # color choice
            )

            # Tooltip
            fig.update_traces(
                hovertemplate=(
                    "<b>WWR: %{x:.2}</b><br>"
                    "<b>EUI: %{customdata[0]:,.0f} kWh</b><br>"
                    "<b>%{y:.1f}% Saving</b><extra></extra>"
                ),
                customdata=df_melt.iloc[1:][["Energy_Outcome(KWH)"]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)
            )

            # Legend at bottom & rename
            fig.for_each_trace(lambda t: t.update(name=t.name.replace("WWR", "WWR")))
            fig.update_layout(
                legend=dict(
                    orientation="h",
                    yanchor="top",
                    y=-0.21,
                    xanchor="center",
                    x=0.5,
                    title=None
                ),
                yaxis_title="% Saving",
                xaxis_title=x_label,
                yaxis=dict(tickformat=".0f"),
                bargap=0.35
            )

            return fig
        
        def make_savings_wwr_barplot_matplotlib(df, title="Energy Savings vs WWR", x_label="Window-to-Wall Ratio"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Drop baseline (first row) to match Plotly
            df_plot = df.iloc[1:].reset_index(drop=True)

            # X positions
            x_values = df_plot["WWR"].values
            x = np.arange(len(x_values))
            width = 0.6

            fig, ax = plt.subplots(figsize=(10,6))

            # Bars
            ax.bar(x, df_plot["PercentSaving"], width, color="blue")

            # Add text labels
            for i, val in enumerate(df_plot["PercentSaving"]):
                ax.text(x[i], val + 0.5, f"{val:.1f}%", ha="center", va="bottom", fontsize=10)

            # X-axis labels
            ax.set_xticks(x)
            ax.set_xticklabels([f"{v:.2f}" for v in x_values])

            # Labels & title
            ax.set_xlabel(x_label)
            ax.set_ylabel("% Saving")
            ax.set_title(title)

            # Grid & layout
            ax.grid(axis='y', alpha=0.3)
            plt.tight_layout()

            return fig

        def make_combined_barplot(df, title, x_label):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            # Reference values for "As Designed"
            reference = df[df["CaseType"] == "As Designed"].iloc[0]
            reference_x = reference["X"]
            reference_energy = reference["Energy_Outcome(KWH)"]

            # % change in Power Density (X-axis)
            df["% Change in Power Density"] = (((reference_x - df["X"]).abs() / reference_x) * 100).round(0)

            # % saving in Energy (Y-axis)
            df["% Saving in Energy"] = (((reference_energy - df["Energy_Outcome(KWH)"]) / reference_energy) * 100).round(0)

            # Bar chart instead of scatter
            fig = px.bar(
                df,
                x="X",
                y="Energy_Outcome(KWH)",
                color="Variable",  # Lighting / Equipment
                pattern_shape="CaseType",  # Different fill patterns for As Designed vs ECM
                barmode="group",
                color_discrete_map={"Lighting": "lightblue", "Equipment": "orange"}
            )

            # Hover with both % saving and % change in PD
            fig.update_traces(
                hovertemplate=(
                    f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                    "<b>Energy Use (kWh)</b>: %{y:.2f}<br>"
                    "<b>% Saving in Energy</b>: %{customdata[0]:.0f}%<br>"
                    "<b>% Change in Power Density</b>: %{customdata[1]:.0f}%"
                    "<extra></extra>"
                ),
                customdata=df[["% Saving in Energy", "% Change in Power Density"]].values
            )

            # Layout
            fig.update_layout(
                # title=title,
                xaxis_title=x_label,
                yaxis_title="Energy Use (kWh)",
                yaxis=dict(tickformat=","),
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=-0.3,
                    xanchor="center",
                    x=0.5
                ),
                bargap=0.25
            )

            return fig

        import plotly.express as px

        def make_savings_gainsWall(df, title="Energy Savings vs R-Value", x_label="Incremental R-Value over As Designed"):
            df = df.copy()
            # -------------------------------------------------
            # 1. CaseType: As Designed per Variable
            # -------------------------------------------------
            df["CaseType"] = "ECM"
            df.loc[df.groupby("Variable").head(1).index, "CaseType"] = "As Designed"

            # -------------------------------------------------
            # 2. Select correct Conduction column per Variable
            # -------------------------------------------------
            df["Conduction"] = np.where(
                df["Variable"].str.contains("wall", case=False),
                df["WALL CONDUCTION"],
                df["ROOF CONDUCTION"]
            )

            # -------------------------------------------------
            # 3. Baseline conduction per Variable
            # -------------------------------------------------
            baseline_map = (
                df[df["CaseType"] == "As Designed"]
                .set_index("Variable")["Conduction"]
                .to_dict()
            )

            # -------------------------------------------------
            # 4. % Saving based on conduction reduction
            # -------------------------------------------------
            df["PercentSaving"] = df.apply(
                lambda r: (
                    (baseline_map[r["Variable"]] - r["Conduction"])
                    / baseline_map[r["Variable"]] * 100
                ) if r["Variable"] in baseline_map else np.nan,
                axis=1
            ).round(1)

            # -------------------------------------------------
            # 5. Select correct R-Value per Variable
            # -------------------------------------------------
            df["R_Value"] = np.where(
                df["Variable"].str.contains("wall", case=False),
                df["R-Value-Wall"],
                df["R-Value-Roof"]
            )

            # -------------------------------------------------
            # 6. Plot ONLY ECM cases
            # -------------------------------------------------
            plot_df = df[df["CaseType"] == "ECM"]

            fig = px.bar(
                plot_df,
                x="R_Value",
                y="PercentSaving",
                color="Variable",
                barmode="group",
                text="PercentSaving",
                labels={
                    "R_Value": x_label,
                    "PercentSaving": "% Reduction in Conduction",
                    "Variable": "Component"
                },
                color_discrete_map={
                    "Wall": "blue",
                    "Roof": "red"
                }
            )

            # -------------------------------------------------
            # 7. Hover info
            # -------------------------------------------------
            fig.update_traces(
                hovertemplate=(
                    "<b>Incremental R-Value: %{x}</b><br>"
                    "<b>Conduction: %{customdata[0]:,.1f}</b><br>"
                    "<b>% Reduction: %{y:.1f}%</b><extra></extra>"
                ),
                customdata=plot_df[["Conduction"]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)
            )

            # -------------------------------------------------
            # 8. Layout
            # -------------------------------------------------
            fig.update_layout(
                title=title,
                yaxis_title="% Reduction in Conduction",
                xaxis_title=x_label,
                yaxis=dict(tickformat=".0f"),
                bargap=0.35,
                legend=dict(
                    orientation="h",
                    yanchor="top",
                    y=-0.22,
                    xanchor="center",
                    x=0.5,
                    title=None
                )
            )

            # -------------------------------------------------
            # 9. Custom X-axis labels (+R)
            # -------------------------------------------------
            unique_rvalues = sorted(plot_df["R_Value"].unique())
            fig.update_xaxes(
                tickvals=unique_rvalues,
                ticktext=[f"+R{v}" for v in unique_rvalues]
            )

            return fig
        
        def make_savings_gainsWindow(df, title="Window Gains vs U-Value", x_label="U-Value (BTU/hr·ft²·F)"):
            df = df.copy()
            # -------------------------------------------------
            # 1. CaseType: As Designed per Variable
            # -------------------------------------------------
            df["CaseType"] = "ECM"
            df.loc[df.groupby("Variable").head(1).index, "CaseType"] = "As Designed"
            # -------------------------------------------------
            # 2. Window conduction column
            # -------------------------------------------------
            conduction_col = "WindowConduction"
            u_col = "BUILDING-Window-U-Value(BTU/HR-SQFT-F)"
            # -------------------------------------------------
            # 3. Baseline conduction per Variable (Window)
            # -------------------------------------------------
            baseline_map = (
                df[df["CaseType"] == "As Designed"]
                .set_index("Variable")[conduction_col]
                .to_dict()
            )

            # -------------------------------------------------
            # 4. % Reduction in Window Conduction
            # -------------------------------------------------
            df["PercentSaving"] = df.apply(
                lambda r: (
                    (baseline_map[r["Variable"]] - r[conduction_col])
                    / baseline_map[r["Variable"]] * 100
                ) if r["Variable"] in baseline_map else np.nan,
                axis=1
            ).round(1)

            # -------------------------------------------------
            # 5. Plot ONLY ECM cases
            # -------------------------------------------------
            plot_df = df[df["CaseType"] == "ECM"]

            fig = px.bar(
                plot_df,
                x=u_col,
                y="PercentSaving",
                text="PercentSaving",
                labels={
                    u_col: x_label,
                    "PercentSaving": "% Reduction in Window Conduction"
                }
            )

            # -------------------------------------------------
            # 6. Hover info
            # -------------------------------------------------
            fig.update_traces(
                marker_color="blue",
                hovertemplate=(
                    "<b>U-Value</b>: %{x:.3f}<br>"
                    "<b>Window Conduction</b>: %{customdata[0]:,.1f}<br>"
                    "<b>% Reduction</b>: %{y:.1f}%<extra></extra>"
                ),
                customdata=plot_df[[conduction_col]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)
            )

            # -------------------------------------------------
            # 7. Layout
            # -------------------------------------------------
            fig.update_layout(
                title=title,
                xaxis_title=x_label,
                yaxis_title="% Reduction in Window Conduction",
                yaxis=dict(tickformat=".0f"),
                bargap=0.35,
                showlegend=False
            )
            return fig
            
        def make_combined_savings_barplot(df, title="Energy Savings vs R-Value", x_label="R-Value (ft²·F·hr/BTU)"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to Baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Melt Wall & Roof into one column for plotting
            df_melt = pd.melt(
                df,
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["R-Value-Wall", "R-Value-Roof"],
                var_name="Variable",
                value_name="R_Value"
            )

            # Plot
            fig = px.bar(
                df_melt.iloc[1:],  # exclude baseline
                x="R_Value",
                y="PercentSaving",
                color="Variable",
                barmode="group",
                text="PercentSaving",
                labels={
                    "R_Value": x_label,
                    "PercentSaving": "% Saving",
                    "Variable": "Component"
                },
                color_discrete_map={"R-Value-Wall": "blue", "R-Value-Roof": "red"}
            )

            # Tooltip
            fig.update_traces(
                hovertemplate=(
                    "<b>R-Value: %{x}</b><br>"
                    "<b>EUI: %{customdata[0]:,.0f} kWh</b><br>"
                    "<b>%{y:.1f}% Saving</b><extra></extra>"
                ),
                customdata=df_melt.iloc[1:][["Energy_Outcome(KWH)"]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)  # Increase text size for % Saving
            )

            # Layout adjustments
            fig.update_layout(
                yaxis_title="% Saving",
                xaxis_title=x_label,
                yaxis=dict(tickformat=".0f"),
                bargap=0.35  # Decrease gap → thicker bars
            )

            # Rename legend items (Wall, Roof)
            fig.for_each_trace(lambda t: t.update(name=t.name.replace("R-Value-", "")))

            fig.update_layout(
                legend=dict(
                    orientation="h",          # horizontal legend
                    yanchor="top",
                    y=-0.21,                  # move below chart
                    xanchor="center",
                    x=0.5,
                    title=None                # remove "Component" title
                )
            )

            # --- Custom x-axis labels ---
            # Get unique R-Values used in plotting
            unique_rvalues = sorted(df_melt["R_Value"].unique())
            # Create ticktext like +R2.5, +R5, ...
            ticktext = [f"+R{v}" for v in unique_rvalues]

            fig.update_xaxes(
                tickvals=unique_rvalues,
                ticktext=ticktext
            )

            return fig
            
        def make_combined_savings_barplot(df, title="Energy Savings vs R-Value", x_label="R-Value (ft²·F·hr/BTU)"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to Baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Melt Wall & Roof into one column for plotting
            df_melt = pd.melt(
                df,
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["R-Value-Wall", "R-Value-Roof"],
                var_name="Variable",
                value_name="R_Value"
            )

            # Plot
            fig = px.bar(
                df_melt.iloc[1:],  # exclude baseline
                x="R_Value",
                y="PercentSaving",
                color="Variable",
                barmode="group",
                text="PercentSaving",
                labels={
                    "R_Value": x_label,
                    "PercentSaving": "% Saving",
                    "Variable": "Component"
                },
                color_discrete_map={"R-Value-Wall": "blue", "R-Value-Roof": "red"}
            )

            # Tooltip
            fig.update_traces(
                hovertemplate=(
                    "<b>R-Value: %{x}</b><br>"
                    "<b>EUI: %{customdata[0]:,.0f} kWh</b><br>"
                    "<b>%{y:.1f}% Saving</b><extra></extra>"
                ),
                customdata=df_melt.iloc[1:][["Energy_Outcome(KWH)"]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)  # Increase text size for % Saving
            )

            # Layout adjustments
            fig.update_layout(
                yaxis_title="% Saving",
                xaxis_title=x_label,
                yaxis=dict(tickformat=".0f"),
                bargap=0.35  # Decrease gap → thicker bars
            )

            # Rename legend items (Wall, Roof)
            fig.for_each_trace(lambda t: t.update(name=t.name.replace("R-Value-", "")))

            fig.update_layout(
                legend=dict(
                    orientation="h",          # horizontal legend
                    yanchor="top",
                    y=-0.21,                  # move below chart
                    xanchor="center",
                    x=0.5,
                    title=None                # remove "Component" title
                )
            )

            # --- Custom x-axis labels ---
            # Get unique R-Values used in plotting
            unique_rvalues = sorted(df_melt["R_Value"].unique())
            # Create ticktext like +R2.5, +R5, ...
            ticktext = [f"+R{v}" for v in unique_rvalues]

            fig.update_xaxes(
                tickvals=unique_rvalues,
                ticktext=ticktext
            )

            return fig

        # Example mapping for LPD/EPD
        lpd_map = {1: 0.1, 2: 0.15, 3: 0.2, 4: 0.25, 5: 0.3}
        def make_savings_lpd_epd_barplot(df, title="EUI Savings vs LPD/EPD", x_label="Power Density (W/ft²)"):
            df = df.copy()

            # Baseline Energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to Baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Extract LPD & EPD codes from FileName (7th and 8th values after splitting by '_')
            df["LPD_Code"] = df["FileName"].str.split("_").str[4].astype(int)
            df["EPD_Code"] = df["FileName"].str.split("_").str[6].astype(int)

            # Map to actual LPD/EPD values
            df["LPD_Value"] = df["LPD_Code"].map(lpd_map)
            df["EPD_Value"] = df["EPD_Code"].map(lpd_map)

            # ✅ Multiply LPD & EPD by 100
            df["LPD_Value"] = df["LPD_Value"] * 100
            df["EPD_Value"] = df["EPD_Value"] * 100

            # Melt into one column for plotting
            df_melt = pd.melt(
                df,
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["LPD_Value", "EPD_Value"],
                var_name="Variable",
                value_name="X_Value"
            )

            # Plot
            fig = px.bar(
                df_melt.iloc[1:],  # exclude baseline
                x="X_Value",
                y="PercentSaving",
                color="Variable",
                barmode="group",
                text="PercentSaving",
                labels={
                    "X_Value": x_label,
                    "PercentSaving": "% Saving",
                    "Variable": "Component"
                },
                # title=title,
                color_discrete_map={"LPD_Value": "blue", "EPD_Value": "red"}
            )

            # Tooltip + Text Styling
            fig.update_traces(
                hovertemplate=(
                    # "<b>Power Density Change(): %{x*100}</b><br>"
                    "<b>EUI: %{customdata[0]:,.0f} kWh</b><br>"
                    "<b>%{y:.1f}% Saving</b><extra></extra>"
                ),
                customdata=df_melt.iloc[1:][["Energy_Outcome(KWH)"]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)
            )

            # Layout → thinner bars
            fig.update_layout(
                yaxis_title="% Saving",
                xaxis_title=x_label,
                yaxis=dict(tickformat=".0f"),
                bargap=0.35,
                legend_title="Component"
            )

            # --- Add this after fig creation ---
            fig.for_each_trace(lambda t: t.update(name=t.name.replace("_Value", "").replace("LPD", "Lighting").replace("EPD", "Equipment")))
            fig.update_layout(
                legend=dict(
                    orientation="h",          # horizontal legend
                    yanchor="top",
                    y=-0.21,                  # move below chart
                    xanchor="center",
                    x=0.5,
                    title=None                # remove "Component" title
                )
            )
            fig.update_xaxes(
                tickvals=[x for x in sorted(df_melt["X_Value"].dropna().unique())],
                ticktext=[f"{int(x)}%" for x in sorted(df_melt["X_Value"].dropna().unique())]
            )

            return fig
        
        def make_savings_lpd_epd_barplot_matplotlib(df, title="EUI Savings vs PD", x_label="Power Density (W/ft²)"):
            df = df.copy()

            # Baseline
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # % Saving
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Extract codes
            df["LPD_Code"] = df["FileName"].str.split("_").str[4].astype(int)
            df["EPD_Code"] = df["FileName"].str.split("_").str[6].astype(int)

            # Map to actual values
            df["LPD_Value"] = df["LPD_Code"].map(lpd_map)
            df["EPD_Value"] = df["EPD_Code"].map(lpd_map)

            # ---- MATCH PLOTLY EXACTLY: drop baseline BEFORE melt ----
            df_plot = df.iloc[1:].copy()   # Equivalent to df_melt.iloc[1:]

            # Melt
            df_melt = df_plot.melt(
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["LPD_Value", "EPD_Value"],
                var_name="Variable",
                value_name="X_Value"
            )

            # Replace names like Plotly
            df_melt["Variable"] = df_melt["Variable"].str.replace("_Value", "", regex=False)
            df_melt["Variable"] = df_melt["Variable"].replace({"LPD": "Lighting", "EPD": "Equipment"})

            # === Build Plot ===

            fig, ax = plt.subplots(figsize=(10, 6))

            # Unique x positions (LPD + EPD per case)
            x_values = sorted(df_melt["X_Value"].unique())
            x = np.arange(len(x_values))  # base positions
            width = 0.35

            # LPD subset
            lighting = df_melt[df_melt["Variable"] == "Lighting"]
            equipment = df_melt[df_melt["Variable"] == "Equipment"]

            # IMPORTANT: Align by x_values to avoid mismatch
            lighting_vals = []
            equipment_vals = []
            for xv in x_values:
                # find the row with this x_value (LPD or EPD)
                l_val = lighting[lighting["X_Value"] == xv]["PercentSaving"]
                e_val = equipment[equipment["X_Value"] == xv]["PercentSaving"]

                lighting_vals.append(l_val.iloc[0] if not l_val.empty else np.nan)
                equipment_vals.append(e_val.iloc[0] if not e_val.empty else np.nan)

            # Plot bars
            ax.bar(x - width/2, lighting_vals, width, label="Lighting", color="blue")
            ax.bar(x + width/2, equipment_vals, width, label="Equipment", color="red")

            # Labels
            ax.set_xlabel(x_label)
            ax.set_ylabel("% Saving")
            ax.set_xticks(x)
            ax.set_xticklabels([str(v) for v in x_values])
            ax.set_title(title)

            # Add text on bars
            for i, val in enumerate(lighting_vals):
                if not np.isnan(val):
                    ax.text(i - width/2, val, f"{val:.1f}%", ha="center", va="bottom", fontsize=10)

            for i, val in enumerate(equipment_vals):
                if not np.isnan(val):
                    ax.text(i + width/2, val, f"{val:.1f}%", ha="center", va="bottom", fontsize=10)

            # Legend (horizontal, bottom, center)
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.15), ncol=2, frameon=False)

            fig.tight_layout()
            return fig

        import plotly.graph_objects as go
        import numpy as np

        def make_combined_plot_gainsWall(df, title, x_label):
            df = df.copy()

            # ---------------------------------
            # 1. Correct CaseType assignment
            # ---------------------------------
            df["CaseType"] = "ECM"
            df.loc[df.groupby("Variable").head(1).index, "CaseType"] = "As Designed"

            # ---------------------------------
            # 2. Select correct Y column
            # ---------------------------------
            df["EnergyGain"] = np.where(
                df["Variable"].str.contains("wall", case=False),
                df["WALL CONDUCTION"],
                np.where(
                    df["Variable"].str.contains("roof", case=False),
                    df["ROOF CONDUCTION"],
                    np.nan
                )
            )

            # ---------------------------------
            # 3. Create figure
            # ---------------------------------
            fig = go.Figure()

            colors = {
                "Wall": "red",
                "Roof": "blue"
            }

            # ---------------------------------
            # 4. Plot ECM
            # ---------------------------------
            ecm_df = df[df["CaseType"] == "ECM"]

            for var in ecm_df["Variable"].unique():
                sub = ecm_df[ecm_df["Variable"] == var]
                color = colors["Wall"] if "wall" in var.lower() else colors["Roof"]

                fig.add_trace(go.Scatter(
                    x=sub["X"],
                    y=sub["EnergyGain"],
                    mode="markers",
                    name=f"{var} - ECM",
                    marker=dict(
                        size=11,
                        color=color,
                        line=dict(width=1, color="black")
                    )
                ))

            # ---------------------------------
            # 5. Plot As Designed (NOW BOTH WALL & ROOF)
            # ---------------------------------
            ad_df = df[df["CaseType"] == "As Designed"]

            for _, row in ad_df.iterrows():
                color = colors["Wall"] if "wall" in row["Variable"].lower() else colors["Roof"]

                fig.add_trace(go.Scatter(
                    x=[row["X"]],
                    y=[row["EnergyGain"]],
                    mode="markers",
                    name=f"{row['Variable']} - As Designed",
                    marker=dict(
                        size=14,
                        color="white",
                        line=dict(width=3, color=color)
                    )
                ))

            # ---------------------------------
            # 6. Layout
            # ---------------------------------
            fig.update_layout(
                title=title,
                xaxis_title=x_label,
                yaxis_title="Conduction Gains (kW)",
                yaxis=dict(tickformat="~s"),
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=-0.35,
                    xanchor="center",
                    x=0.5
                ),
                showlegend=True
            )

            return fig

        def make_combined_plot(df, title, x_label):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            reference_x = df[df["CaseType"] == "As Designed"].iloc[0]["X"]
            df["PercentChange"] = (((reference_x - df["X"]).abs() / reference_x) * 100).round(0)

            # Human readable numbers
            def human_readable(num):
                if num >= 1_000_000_000:
                    return f"{num/1_000_000_000:.1f}B (Billion)"
                elif num >= 1_000_000:
                    return f"{num/1_000_000:.1f}M (Million)"
                elif num >= 1_000:
                    return f"{num/1_000:.1f}k (Thousand)"
                else:
                    return str(num)

            df["Energy_human"] = df["Energy_Outcome(KWH)"].apply(human_readable)
            df["PercentChangeAdj"] = df["PercentChange"] + 5

            # Assign only red and blue alternately to each parameter
            colors = ["red", "blue"]
            unique_vars = df["Variable"].unique()
            color_map = {var: colors[i % 2] for i, var in enumerate(unique_vars)}

            fig = go.Figure()

            for variable in unique_vars:
                color = color_map[variable]
                subset = df[df["Variable"] == variable]

                # ECM – filled circle
                ecm_df = subset[subset["CaseType"] == "ECM"]
                fig.add_trace(go.Scatter(
                    x=ecm_df["X"],
                    y=ecm_df["Energy_Outcome(KWH)"],
                    mode="markers",
                    name=f"{variable} - ECM",
                    marker=dict(size=12, color=color, line=dict(width=1, color="black"), symbol="circle", opacity=0.9),
                    customdata=ecm_df[["Energy_human", "PercentChangeAdj"]].values,
                    hovertemplate=(
                        f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                        "<b>Energy Use</b>: %{customdata[0]}<br>"
                        # "<b>% Change from As Designed</b>: %{customdata[1]:.0f}%<extra></extra>"
                    )
                ))

                # As Designed – unfilled circle
                ad_df = subset[subset["CaseType"] == "As Designed"]
                fig.add_trace(go.Scatter(
                    x=ad_df["X"],
                    y=ad_df["Energy_Outcome(KWH)"],
                    mode="markers",
                    name=f"{variable} - As Designed",
                    marker=dict(size=12, color="white", line=dict(width=2, color=color), symbol="circle", opacity=0.95),
                    customdata=ad_df[["Energy_human", "PercentChangeAdj"]].values,
                    hovertemplate=(
                        f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                        "<b>Energy Use</b>: %{customdata[0]}<br>"
                        # "<b>% Change from As Designed</b>: %{customdata[1]:.0f}%<extra></extra>"
                    )
                ))

            fig.update_layout(
                # title=title,
                xaxis_title=x_label,
                yaxis_title="Energy Use (kWh, K=Thousand, M=Million, B=Billion)",
                yaxis=dict(tickformat="~s"),
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=-0.3,
                    xanchor="center",
                    x=0.5
                ),
                showlegend=True
            )

            return fig
        
        from matplotlib.backends.backend_pdf import PdfPages

        def make_combined_plot_matplotlib(df, title, x_label):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            reference_x = df[df["CaseType"] == "As Designed"].iloc[0]["X"]
            df["PercentChange"] = (((reference_x - df["X"]).abs() / reference_x) * 100).round(0)

            def human_readable(num):
                if num >= 1_000_000_000:
                    return f"{num/1_000_000_000:.1f}B"
                elif num >= 1_000_000:
                    return f"{num/1_000_000:.1f}M"
                elif num >= 1_000:
                    return f"{num/1_000:.1f}k"
                else:
                    return str(num)

            df["Energy_human"] = df["Energy_Outcome(KWH)"].apply(human_readable)
            df["PercentChangeAdj"] = df["PercentChange"] + 5

            colors = ["red", "blue"]
            unique_vars = df["Variable"].unique()
            color_map = {var: colors[i % 2] for i, var in enumerate(unique_vars)}

            fig, ax = plt.subplots(figsize=(10, 6))

            for variable in unique_vars:
                color = color_map[variable]
                subset = df[df["Variable"] == variable]

                # ECM – filled circle
                ecm_df = subset[subset["CaseType"] == "ECM"]
                ax.scatter(
                    ecm_df["X"],
                    ecm_df["Energy_Outcome(KWH)"],
                    s=100,
                    color=color,
                    edgecolor="black",
                    alpha=0.9,
                    label=f"{variable} - ECM",
                )

                # As Designed – unfilled circle
                ad_df = subset[subset["CaseType"] == "As Designed"]
                ax.scatter(
                    ad_df["X"],
                    ad_df["Energy_Outcome(KWH)"],
                    s=100,
                    facecolors="none",
                    edgecolor=color,
                    linewidth=2,
                    alpha=0.95,
                    label=f"{variable} - As Designed",
                )

            ax.set_xlabel(x_label, fontsize=12)
            ax.set_ylabel("Energy Use (kWh, K=Thousand, M=Million, B=Billion)", fontsize=12)
            ax.legend(loc="lower center", bbox_to_anchor=(0.5, -0.3), ncol=2)
            ax.grid(alpha=0.3)

            plt.tight_layout()
            return fig
        
        def make_savings_barplot_matplotlib(df, title="Energy Savings vs R-Value", x_label="R-Value (ft²·F·hr/BTU)"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Variables to plot
            variables = ["R-Value-Wall", "R-Value-Roof"]
            colors = ["blue", "red"]

            # Exclude baseline row for plotting
            plot_df = df.iloc[1:].reset_index(drop=True)

            # Bar positions
            x = np.arange(len(plot_df))
            width = 0.35  # width of each bar

            fig, ax = plt.subplots(figsize=(10,6))

            # Plot each variable
            for i, var in enumerate(variables):
                ax.bar(
                    x + i*width,
                    plot_df["PercentSaving"],
                    width=width,
                    color=colors[i],
                    label=var.replace("R-Value-", "")
                )

                # Add value labels
                for xi, val in zip(x + i*width, plot_df["PercentSaving"]):
                    ax.text(xi, val + 0.5, f"{val:.1f}%", ha='center', va='bottom', fontsize=12)

            # Set x-ticks to be in the center of the group
            ax.set_xticks(x + width/2)
            # Use actual R-values as tick labels (for Wall)
            ax.set_xticklabels([f"{val}" for val in plot_df["R-Value-Wall"]], fontsize=12)

            ax.set_xlabel(x_label, fontsize=12)
            ax.set_ylabel("% Saving", fontsize=12)
            ax.set_title(title, fontsize=14)

            # Legend below chart
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.15), ncol=2)
            ax.grid(alpha=0.3)

            plt.tight_layout()
            return fig
        
        def make_scatter_case_type_matplotlib(df, col, x_label, y_min=None, y_max=None, title="Energy Use"):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            fig, ax = plt.subplots(figsize=(8,6))

            colors = {"As Designed": "red", "ECM": "blue"}
            marker_size = 120  # Matplotlib marker size in points² (roughly similar to Plotly size=12)
            alpha = 0.85

            for case_type, group in df.groupby("CaseType"):
                ax.scatter(
                    group[col],
                    group["Energy_Outcome(KWH)"],
                    s=marker_size,
                    c=colors[case_type],
                    alpha=alpha,
                    label=case_type,
                    edgecolor="black"
                )
            
            ax.set_xlabel(x_label)
            ax.set_ylabel("Energy Use (kWh, K=Thousand, M=Million, B=Billion)")
            ax.set_title(title)
            ax.legend(loc='upper center', bbox_to_anchor=(0.5, -0.15), ncol=2, frameon=False)
            if y_min is not None and y_max is not None:
                ax.set_ylim([y_min, y_max])
            ax.grid(alpha=0.3)
            plt.tight_layout()
            return fig
        
        def make_savings_sc_barplot_matplotlib(df, title="Energy Savings vs SC", x_label="Shading Coefficient (SC)"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to Baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Melt EUI & WWR into one column for plotting (if multiple types, else just use WWR)
            df_melt = pd.melt(
                df,
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["SC"],  # replace/add if multiple types
                var_name="Variable",
                value_name="WWR_Value"
            )

            # Plot
            fig = px.bar(
                df_melt.iloc[1:],  # exclude baseline
                x="WWR_Value",
                y="PercentSaving",
                color="Variable",
                barmode="group",
                text="PercentSaving",
                labels={
                    "WWR_Value": x_label,
                    "PercentSaving": "% Saving",
                    "Variable": "Component"
                },
                color_discrete_map={"SC": "blue"}  # color choice
            )

            # Tooltip
            fig.update_traces(
                hovertemplate=(
                    "<b>SC: %{x:.2}</b><br>"
                    "<b>EUI: %{customdata[0]:,.0f} kWh</b><br>"
                    "<b>%{y:.1f}% Saving</b><extra></extra>"
                ),
                customdata=df_melt.iloc[1:][["Energy_Outcome(KWH)"]].to_numpy(),
                textposition="outside",
                textfont=dict(size=16)
            )

            # Legend at bottom & rename
            fig.for_each_trace(lambda t: t.update(name=t.name.replace("SC", "SC")))
            fig.update_layout(
                legend=dict(
                    orientation="h",
                    yanchor="top",
                    y=-0.21,
                    xanchor="center",
                    x=0.5,
                    title=None
                ),
                yaxis_title="% Saving",
                xaxis_title=x_label,
                yaxis=dict(tickformat=".0f"),
                bargap=0.35
            )

            return fig
        
        import plotly.express as px
        def make_combined_plot_window_matplotlib(df, title, x_label):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            reference_x = df[df["CaseType"] == "As Designed"].iloc[0]["X"]
            df["PercentChange"] = (((reference_x - df["X"]).abs() / reference_x) * 100).round(0)

            # Human readable numbers
            def human_readable(num):
                if num >= 1_000_000_000:
                    return f"{num/1_000_000_000:.1f}B"
                elif num >= 1_000_000:
                    return f"{num/1_000_000:.1f}M"
                elif num >= 1_000:
                    return f"{num/1_000:.1f}k"
                else:
                    return str(num)

            df["Energy_human"] = df["Energy_Outcome(KWH)"].apply(human_readable)

            # Normalize U-value for marker size
            u_min = df["BUILDING-Window-U-Value(BTU/HR-SQFT-F)"].min()
            u_max = df["BUILDING-Window-U-Value(BTU/HR-SQFT-F)"].max()
            size_min = 50
            size_max = 300

            def u_to_size(u):
                return size_min + (u - u_min) / (u_max - u_min) * (size_max - size_min)

            # Create figure and axis
            fig, ax = plt.subplots(figsize=(10, 6))
            
            unique_vars = df["Variable"].unique()
            colors = ["red", "blue"]
            color_map = {var: colors[i % 2] for i, var in enumerate(unique_vars)}

            for variable in unique_vars:
                subset = df[df["Variable"] == variable]
                
                # ECM – filled circle
                ecm_df = subset[subset["CaseType"] == "ECM"]
                ax.scatter(
                    ecm_df["X"],
                    ecm_df["Energy_Outcome(KWH)"],
                    s=ecm_df["BUILDING-Window-U-Value(BTU/HR-SQFT-F)"].apply(u_to_size),
                    c="blue",
                    label=f"ECM",
                    alpha=0.7,
                    edgecolors="black"
                )
                
                # As Designed – unfilled circle
                ad_df = subset[subset["CaseType"] == "As Designed"]
                ax.scatter(
                    ad_df["X"],
                    ad_df["Energy_Outcome(KWH)"],
                    s=ad_df["BUILDING-Window-U-Value(BTU/HR-SQFT-F)"].apply(u_to_size),
                    facecolors='red',
                    edgecolors="red",
                    linewidths=2,
                    label=f"As Designed"
                )

            ax.set_title(title)
            ax.set_xlabel(x_label)
            ax.set_ylabel("Energy Use (kWh)")
            ax.legend()
            ax.grid(True)

            return fig, ax

        def make_combined_scatter(df, title, x_label):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            reference_x = df[df["CaseType"] == "As Designed"].iloc[0]["X"]
            df["PercentChange"] = (((reference_x - df["X"]).abs() / reference_x) * 100).round(0)

            # Human readable numbers
            def human_readable(num):
                if num >= 1_000_000_000:
                    return f"{num/1_000_000_000:.1f}B (Billion)"
                elif num >= 1_000_000:
                    return f"{num/1_000_000:.1f}M (Million)"
                elif num >= 1_000:
                    return f"{num/1_000:.1f}k (Thousand)"
                else:
                    return str(num)

            df["Energy_human"] = df["Energy_Outcome(KWH)"].apply(human_readable)
            df["PercentChangeAdj"] = df["PercentChange"] + 5

            # Assign only red and blue alternately to each parameter
            colors = ["red", "blue"]
            unique_vars = df["Variable"].unique()
            color_map = {var: colors[i % 2] for i, var in enumerate(unique_vars)}

            # Normalize U-value for marker size
            u_col = "BUILDING-Window-U-Value(BTU/HR-SQFT-F)"
            u_min = df[u_col].min(skipna=True)
            u_max = df[u_col].max(skipna=True)
            size_min = 8
            size_max = 25

            # Warn if NaN U-values found
            if df[u_col].isna().any():
                print("⚠️ Warning: Missing U-values found. Using minimum U-value for those entries.")
                df[u_col] = df[u_col].fillna(u_min)

            def u_to_size(u):
                # Handle identical or NaN min/max cases
                if pd.isna(u_min) or pd.isna(u_max) or u_min == u_max or pd.isna(u):
                    return size_min  # fallback default
                # Scale U-value to marker size
                return size_min + (u - u_min) / (u_max - u_min) * (size_max - size_min)

            # Create Plotly figure
            fig = go.Figure()

            for variable in unique_vars:
                color = color_map[variable]
                subset = df[df["Variable"] == variable]

                # ECM – filled circle
                ecm_df = subset[subset["CaseType"] == "ECM"]
                fig.add_trace(go.Scatter(
                    x=ecm_df["X"],
                    y=ecm_df["Energy_Outcome(KWH)"],
                    mode="markers",
                    name="ECM",
                    marker=dict(
                        size=ecm_df[u_col].apply(u_to_size),
                        color="blue",
                        line=dict(width=1, color="blue"),
                        symbol="circle",
                        opacity=0.9
                    ),
                    customdata=ecm_df[["Energy_human", "PercentChangeAdj", u_col]].values,
                    hovertemplate=(
                        f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                        "<b>Energy Use</b>: %{customdata[0]}<br>"
                        "<b>U-Value</b>: %{customdata[2]:.3f} BTU/hr·ft²·°F<extra></extra>"
                    )
                ))

                # As Designed – unfilled circle
                ad_df = subset[subset["CaseType"] == "As Designed"]
                fig.add_trace(go.Scatter(
                    x=ad_df["X"],
                    y=ad_df["Energy_Outcome(KWH)"],
                    mode="markers",
                    name="As Designed",
                    marker=dict(
                        size=ad_df[u_col].apply(u_to_size),
                        color="red",
                        line=dict(width=2, color="red"),
                        symbol="circle",
                        opacity=0.95
                    ),
                    customdata=ad_df[["Energy_human", "PercentChangeAdj", u_col]].values,
                    hovertemplate=(
                        f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                        "<b>Energy Use</b>: %{customdata[0]}<br>"
                        "<b>U-Value</b>: %{customdata[2]:.3f} BTU/hr·ft²·°F<extra></extra>"
                    )
                ))

            # Layout formatting
            fig.update_layout(
                title=title,
                xaxis_title=x_label,
                yaxis_title="Energy Use (kWh, K=Thousand, M=Million, B=Billion)",
                yaxis=dict(tickformat="~s"),
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=-0.3,
                    xanchor="center",
                    x=0.5
                ),
                showlegend=True
            )

            return fig

        def make_savings_window_matplotlib(df, title="Energy Savings vs SC", x_label="Shading Coefficient (SC)"):
            df = df.copy()

            # Baseline energy
            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]

            # Compute % Saving relative to Baseline
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Exclude baseline row for plotting
            plot_df = df.iloc[1:].copy()

            # Bar plot
            fig, ax = plt.subplots(figsize=(10, 6))
            
            x_values = plot_df["SC"]
            y_values = plot_df["PercentSaving"]

            bars = ax.bar(x_values, y_values, color="blue", edgecolor="black", width=0.6)

            # Annotate bars with percent savings
            for bar, percent in zip(bars, y_values):
                height = bar.get_height()
                ax.text(bar.get_x() + bar.get_width()/2, height + 0.5, f"{percent}%", ha='center', va='bottom', fontsize=12)

            # Labels and title
            ax.set_xlabel(x_label)
            ax.set_ylabel("% Saving")
            ax.set_title(title)
            ax.grid(axis="y", linestyle="--", alpha=0.7)
            return fig, ax
        
        def make_combined_plot_gainsWindow(df, title, x_label):
            df = df.copy()
            df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")

            reference_x = df[df["CaseType"] == "As Designed"].iloc[0]["X"]
            df["PercentChange"] = (((reference_x - df["X"]).abs() / reference_x) * 100).round(0)

            # Human readable numbers
            def human_readable(num):
                if num >= 1_000_000_000:
                    return f"{num/1_000_000_000:.1f}B (Billion)"
                elif num >= 1_000_000:
                    return f"{num/1_000_000:.1f}M (Million)"
                elif num >= 1_000:
                    return f"{num/1_000:.1f}k (Thousand)"
                else:
                    return str(num)

            df["Energy_human"] = df["WindowConduction"].apply(human_readable)
            df["PercentChangeAdj"] = df["PercentChange"] + 5

            # Assign only red and blue alternately to each parameter
            colors = ["red", "blue"]
            unique_vars = df["Variable"].unique()
            color_map = {var: colors[i % 2] for i, var in enumerate(unique_vars)}

            # Normalize U-value for marker size
            u_col = "BUILDING-Window-U-Value(BTU/HR-SQFT-F)"
            u_min = df[u_col].min(skipna=True)
            u_max = df[u_col].max(skipna=True)
            size_min = 8
            size_max = 25

            # Warn if NaN U-values found
            if df[u_col].isna().any():
                print("⚠️ Warning: Missing U-values found. Using minimum U-value for those entries.")
                df[u_col] = df[u_col].fillna(u_min)

            def u_to_size(u):
                # Handle identical or NaN min/max cases
                if pd.isna(u_min) or pd.isna(u_max) or u_min == u_max or pd.isna(u):
                    return size_min  # fallback default
                # Scale U-value to marker size
                return size_min + (u - u_min) / (u_max - u_min) * (size_max - size_min)

            # Create Plotly figure
            fig = go.Figure()

            for variable in unique_vars:
                color = color_map[variable]
                subset = df[df["Variable"] == variable]

                # ECM – filled circle
                ecm_df = subset[subset["CaseType"] == "ECM"]
                fig.add_trace(go.Scatter(
                    x=ecm_df["X"],
                    y=ecm_df["WindowConduction"],
                    mode="markers",
                    name="ECM",
                    marker=dict(
                        size=ecm_df[u_col].apply(u_to_size),
                        color="blue",
                        line=dict(width=1, color="blue"),
                        symbol="circle",
                        opacity=0.9
                    ),
                    customdata=ecm_df[["Energy_human", "PercentChangeAdj", u_col]].values,
                    hovertemplate=(
                        f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                        "<b>Window Conduction</b>: %{customdata[0]}<br>"
                        "<b>U-Value</b>: %{customdata[2]:.3f} BTU/hr·ft²·°F<extra></extra>"
                    )
                ))

                # As Designed – unfilled circle
                ad_df = subset[subset["CaseType"] == "As Designed"]
                fig.add_trace(go.Scatter(
                    x=ad_df["X"],
                    y=ad_df["WindowConduction"],
                    mode="markers",
                    name="As Designed",
                    marker=dict(
                        size=ad_df[u_col].apply(u_to_size),
                        color="red",
                        line=dict(width=2, color="red"),
                        symbol="circle",
                        opacity=0.95
                    ),
                    customdata=ad_df[["Energy_human", "PercentChangeAdj", u_col]].values,
                    hovertemplate=(
                        f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                        "<b>Window Conduction</b>: %{customdata[0]}<br>"
                        "<b>U-Value</b>: %{customdata[2]:.3f} BTU/hr·ft²·°F<extra></extra>"
                    )
                ))

            # Layout formatting
            fig.update_layout(
                title=title,
                xaxis_title=x_label,
                yaxis_title="Glazing Gains(kW)",
                yaxis=dict(tickformat="~s"),
                legend=dict(
                    orientation="h",
                    yanchor="bottom",
                    y=-0.3,
                    xanchor="center",
                    x=0.5
                ),
                showlegend=True
            )

            return fig
        
        def make_combined_savings_barplot_matplotlib(df, title="Energy Savings vs R-Value", x_label="R-Value"):
            df = df.copy()

            baseline_energy = df.iloc[0]["Energy_Outcome(KWH)"]
            df["PercentSaving"] = ((baseline_energy - df["Energy_Outcome(KWH)"]) / baseline_energy * 100).round(1)

            # Drop baseline
            df_plot = df.iloc[1:].reset_index(drop=True)

            # Melt Wall & Roof like Plotly
            df_melt = df_plot.melt(
                id_vars=["Energy_Outcome(KWH)", "PercentSaving"],
                value_vars=["R-Value-Wall", "R-Value-Roof"],
                var_name="Variable",
                value_name="R_Value"
            )

            df_melt["R_Value"] = df_melt["R_Value"].astype(float)  # ensure numeric

            # X positions for grouped bars
            x = np.arange(len(df_plot))
            width = 0.35

            fig, ax = plt.subplots(figsize=(10,6))

            # Separate Wall & Roof
            wall_data = df_melt[df_melt["Variable"] == "R-Value-Wall"]
            roof_data = df_melt[df_melt["Variable"] == "R-Value-Roof"]

            # Bars
            ax.bar(x - width/2, wall_data["PercentSaving"].values, width, color="blue", label="Wall")
            ax.bar(x + width/2, roof_data["PercentSaving"].values, width, color="red", label="Roof")

            # Add text
            for i in range(len(x)):
                ax.text(x[i] - width/2, wall_data["PercentSaving"].values[i]+0.5,
                        f"{wall_data['PercentSaving'].values[i]:.1f}%", ha="center", va="bottom", fontsize=10)
                ax.text(x[i] + width/2, roof_data["PercentSaving"].values[i]+0.5,
                        f"{roof_data['PercentSaving'].values[i]:.1f}%", ha="center", va="bottom", fontsize=10)

            # X-axis labels
            ax.set_xticks(x)
            ax.set_xticklabels([str(v) for v in wall_data["R_Value"].values])

            # Labels, legend, grid
            ax.set_xlabel(x_label)
            ax.set_ylabel("% Saving")
            ax.set_title(title)
            ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.15), ncol=2, frameon=False)
            ax.grid(axis="y", alpha=0.3)
            plt.tight_layout()

            return fig

        # --- Now update 'cases' for the loop (remove Wall, Roof, Light, Equip) ---
        cases = [
            ("Energy Use vs WWR", wwr_df, "WWR", "Window-to-Wall Ratio"),
            ("Energy Use vs Orientation", orient_df, "Orientation", "Orientation (°)"),
        ]

        # 1️⃣ Generate scatter plots for all cases
        figures_scatter = []
        for exp_title, df, col, x_label in cases:
            if col in df.columns and not df.empty:
                fig = make_scatter_case_type_matplotlib(df, col, x_label, y_min=y_min, y_max=y_max, title=exp_title)
                figures_scatter.append(fig)

        # 2️⃣ Generate your other figures
        figA = make_combined_plot_matplotlib(wall_roof_df, "Energy Use vs Wall & Roof", "U-Value (BTU/hr·ft²·F)")
        figB = make_savings_lpd_epd_barplot_matplotlib(light_equip_df, "Energy Savings vs R-Value", "R-Value (ft²·F·hr/BTU)")
        figC = make_combined_plot_matplotlib(light_equip_df, "Energy Use vs Light & Equipment", "Power Density (W/ft²)")
        figD = make_savings_lpd_epd_barplot_matplotlib(light_equip_df, "Energy Savings vs LPD/EPD", "Power Density (W/ft²)")
        figE = make_savings_wwr_barplot_matplotlib(wwr_df, "Energy Savings vs WWR", "Window-to-Wall Ratio")
        figF, ax = make_combined_plot_window_matplotlib(glazing_df, "Energy Use vs Shading Coefficient(SC)", "Shading Coefficient(SC)")
        figG, ax = make_savings_window_matplotlib(glazing_df, "Energy Saving vs Shading Coefficient(SC)", "Shading Coefficient(SC)")
        
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        # Load existing presentation
        ppt_path = "database/Report_Output.pptx"
        prs = Presentation(ppt_path)
        
        # Your figures
        figures = [figA, figB, figC, figD] + figures_scatter + [figE] + [figF] + [figG]
        slide_indices = [1, 2, 3, 4]  # slides to insert charts
        fig_idx = 0

        for slide_num in slide_indices:
            slide = prs.slides[slide_num]

            # Print all shapes on slide to identify placeholders/textboxes
            # (optional, only for debugging)
            # for i, shape in enumerate(slide.shapes):
            #     print(i, shape.name, shape.shape_type)

            # Example: assume slide has 2 placeholders/textboxes to insert images
            # Adjust indexes based on your PPT
            placeholder_idxs = [1, 2]  # placeholder index where you want left/right chart

            for idx in placeholder_idxs:
                if fig_idx >= len(figures):
                    break

                fig = figures[fig_idx]
                img_stream = io.BytesIO()
                fig.savefig(img_stream, format='PNG', bbox_inches='tight')
                img_stream.seek(0)

                # Use shape coordinates for placement
                shape = slide.shapes[idx]
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
                fig_idx += 1

        # Save to buffer for Streamlit
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        # Common y-axis range
        y_min = combined_Data_expanded["Energy_Outcome(KWH)"].min()
        y_max = combined_Data_expanded["Energy_Outcome(KWH)"].max()

        col1, col2 = st.columns(2)
        with col1:
            fig1 = make_combined_plot(wall_roof_df, "Energy Use vs Wall & Roof", "U-Value (BTU/hr·ft²·F)")
            st.plotly_chart(fig1, use_container_width=True)
        with col2:
            fig2 = make_combined_savings_barplot(wall_roof_df, "Energy Savings vs R-Value", "Incremental R-Value over As Designed")
            st.plotly_chart(fig2, use_container_width=True)
        _, colN2 = st.columns(2)
        with colN2:
            st.markdown("""**Note:**  *X-axis represents the incremental R-Value over the existing construction.*""",unsafe_allow_html=True)

        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_combined_plot(light_equip_df, "Energy Use vs Lighting & Equipment", "Power Density (W/ft²)"), use_container_width=True)
        with col2:
            st.plotly_chart(make_savings_lpd_epd_barplot(light_equip_df, "EUI Savings vs LPD/EPD", "Savings(%) over As Designed Power Density"), use_container_width=True)
        _, colN2 = st.columns(2)
        with colN2:
            st.markdown("""**Note:**  *X-axis represents the Savings (%) over Existing Lighting and Equipment.*""", unsafe_allow_html=True)

        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_combined_scatter(glazing_df, "Energy Use vs Shading Coefficient (SC)", "Shading Coefficient (SC)"), use_container_width=True)
        with col2:
            st.plotly_chart(make_savings_sc_barplot_matplotlib(glazing_df, "EUI Savings vs Shading Coefficient (SC)", "Shading Coefficient (SC)"), use_container_width=True)
        colN1, _ = st.columns(2)
        with colN1:
            st.markdown("""**Note:**  *Size of Points represented by U-Value (BTU/hr·ft²·F).*""", unsafe_allow_html=True)

        # Loop through plots
        for exp_title, df, col, x_label in cases:
            if col in df.columns and not df.empty:
                col1, col2 = st.columns(2)
                with col1:
                    df = df.copy()
                    df["CaseType"] = df.index.map(lambda i: "As Designed" if i == 0 else "ECM")
                    # Scatter plot without CaseType
                    fig = px.scatter(
                        df,
                        x=col,
                        y="Energy_Outcome(KWH)",
                        color="CaseType",
                        color_discrete_map={"As Designed": "red", "ECM": "blue"}
                    )

                    # Make points bigger
                    fig.update_traces(marker=dict(size=12, opacity=0.85))
                    fig.update_layout(legend_title_text="")

                    # Custom hover template
                    fig.update_traces(
                        hovertemplate=(
                            f"<b>{x_label}</b>: %{{x:.2f}}<br>"
                            "<b>Energy Use (kWh)</b>: %{y:.2f}"
                            "<extra></extra>"
                        )
                    )

                    # Layout settings
                    fig.update_layout(
                        xaxis_title=x_label,
                        yaxis_title="Energy Use (kWh, K=Thousand, M=Million, B=Billion)",
                        yaxis=dict(range=[y_min, y_max], tickformat="~s"),
                        legend=dict(
                            orientation="h", 
                            yanchor="bottom",
                            y=-0.3,
                            xanchor="center",
                            x=0.5
                        ),
                        showlegend=True
                    )

                    st.plotly_chart(fig, use_container_width=True)
                
                with col2:
                    st.plotly_chart(
                        make_savings_wwr_barplot(
                            wwr_df,
                            title="Energy Savings vs WWR",
                            x_label="Window-to-Wall Ratio"
                        ),
                        use_container_width=True
                    )
        
        st.markdown(
            """
            <hr style="border: none; height: 2px; background-color: red; margin: 25px 0;">
            <h3 style="color: red; text-align: center; margin-top: -15px;">
                Gains / Losses
            </h3>
            """,
            unsafe_allow_html=True
        )

        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_combined_plot_gainsWall(wall_roof_df, "Conduction Gains vs Wall & Roof", "U-Value (BTU/hr·ft²·F)"), use_container_width=True, key="wall_roof_gains_plot")
        with col2:
            st.plotly_chart(make_savings_gainsWall(wall_roof_df, "Conduction Gains vs Wall & Roof Assemblies", "Incremental R-Value over As Designed"), use_container_width=True)
        _, colN2 = st.columns(2)
        with colN2:
            st.markdown("""**Note:**  *X-axis represents the incremental R-Value over the existing construction.*""",unsafe_allow_html=True)
        glazing_df = glazing_df.assign(WindowConduction=glazing_df["WINDOW GLASS+FRM COND"] + glazing_df["WINDOW GLASS SOLAR"])
        # st.write(glazing_df)
        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_combined_plot_gainsWindow(glazing_df, "Glazing Conduction vs Shading Coefficient", "Shading Coefficient (SC)"), use_container_width=True, key="window_gains_plot")
        with col2:
            st.plotly_chart(make_savings_gainsWindow(glazing_df, "Glazing Gains (%) vs Shading", "Shading Coefficient (SC)"), use_container_width=True, key="gains%")
        colN1, _ = st.columns(2)
        with colN1:
            st.markdown("""**Note:**  *Size of Points represented by U-Value (BTU/hr·ft²·F).*""", unsafe_allow_html=True)
        
        col1, _ = st.columns(2)
        with col1:
            st.download_button(
                label="Download Report",
                data=pptx_buffer,
                file_name="Report_Output.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
        # st.write(combined_Data)
        fixed_cols = [
            "Location", 
            "Floor-Total-Above-Grade(SQFT)", "Floor-Total-Below-Grade(SQFT)",
            "Floor-Total-Conditioned-Grade(SQFT)", "Floor-Total-UnConditioned-Grade(SQFT)",
            "Wall-Total-Above-Grade(SQFT)",
            "N-Wall-Area(SQFT)", "S-Wall-Area(SQFT)", "E-Wall-Area(SQFT)", "W-Wall-Area(SQFT)",
            "SE-Wall-Area(SQFT)", "SW-Wall-Area(SQFT)", "NE-Wall-Area(SQFT)", "NW-Wall-Area(SQFT)",
            "ROOF-AREA(SQFT)", "ALL WALLS-Wall-AREA(SQFT)", "WALLS+ROOFS-AREA(SQFT)",
            "UNDERGRND-Wall-AREA(SQFT)", "BUILDING-Wall-AREA(SQFT)",
            "N-Window-Area(SQFT)", "S-Window-Area(SQFT)", "E-Window-Area(SQFT)", "W-Window-Area(SQFT)",
            "SE-Window-Area(SQFT)", "SW-Window-Area(SQFT)", "NE-Window-Area(SQFT)", "NW-Window-Area(SQFT)",
            "ROOF-Window-Area(SQFT)", "ALL WALLS-Window-AREA(SQFT)", "WALLS+ROOFS-Window-AREA(SQFT)",
            "UNDERGRND-Window-AREA(SQFT)", "BUILDING-Window-AREA(SQFT)"
        ]
        drop_cols = ["ProjectName"]

        # fixed dataframe (only one row is enough since all values are same)
        fixed_df = combined_Data[fixed_cols].iloc[[0]].round(2)
        variable_df = combined_Data.drop(columns=fixed_cols).round(2)
        variable_df = variable_df.drop(columns=drop_cols, axis=1)

        rename_map = {
            "Location": "Location",

            # Floor Areas
            "Floor-Total-Above-Grade(SQFT)": "Above Grade Floor (sqft)",
            "Floor-Total-Below-Grade(SQFT)": "Below Grade Floor (sqft)",
            "Floor-Total-Conditioned-Grade(SQFT)": "Conditioned Floor (sqft)",
            "Floor-Total-UnConditioned-Grade(SQFT)": "Unconditioned Floor (sqft)",

            # Wall Areas
            "Wall-Total-Above-Grade(SQFT)": "Above Grade Wall (sqft)",
            "N-Wall-Area(SQFT)": "North Wall (sqft)",
            "S-Wall-Area(SQFT)": "South Wall (sqft)",
            "E-Wall-Area(SQFT)": "East Wall (sqft)",
            "W-Wall-Area(SQFT)": "West Wall (sqft)",
            "SE-Wall-Area(SQFT)": "SE Wall (sqft)",
            "SW-Wall-Area(SQFT)": "SW Wall (sqft)",
            "ROOF-AREA(SQFT)": "Roof Area (sqft)",
            "ALL WALLS-Wall-AREA(SQFT)": "All Walls Area (sqft)",
            "WALLS+ROOFS-AREA(SQFT)": "Walls+Roof Area (sqft)",
            "UNDERGRND-Wall-AREA(SQFT)": "Underground Wall (sqft)",
            "BUILDING-Wall-AREA(SQFT)": "Building Wall Total (sqft)",

            # Window Areas
            "N-Window-Area(SQFT)": "North Window (sqft)",
            "S-Window-Area(SQFT)": "South Window (sqft)",
            "E-Window-Area(SQFT)": "East Window (sqft)",
            "W-Window-Area(SQFT)": "West Window (sqft)",
            "SE-Window-Area(SQFT)": "SE Window (sqft)",
            "SW-Window-Area(SQFT)": "SW Window (sqft)",
            "ROOF-Window-Area(SQFT)": "Roof Window (sqft)",
            "ALL WALLS-Window-AREA(SQFT)": "All Walls Window (sqft)",
            "WALLS+ROOFS-Window-AREA(SQFT)": "Walls+Roof Window (sqft)",
            "UNDERGRND-Window-AREA(SQFT)": "Underground Window (sqft)",
            "BUILDING-Window-AREA(SQFT)": "Building Window Total (sqft)"
        }

        # Apply renaming
        fixed_df = fixed_df.rename(columns=rename_map)
        variable_df = variable_df.rename(columns=rename_map)

        colA, colB = st.columns(2)
        with colA:
            with st.expander("Building Information"):
                cData = combined_Data.copy()
                combined_Data = combined_Data.drop(combined_Data.columns[8:], axis=1)
                combined_Data = combined_Data.drop(columns=["ProjectName"])
                combined_Data = combined_Data.T
                # combined_Data["FileName"] = combined_Data["FileName"].str.replace(r'_[^_]+(?=\.[^.]+$|$)', f'_{user_nm}', regex=True)
                combined_Data = combined_Data.drop(combined_Data.columns[1:], axis=1)
                fixed_df = fixed_df.reset_index(drop=True)
                numeric_cols = fixed_df.select_dtypes(include='number').columns  # all numeric columns
                fixed_df[numeric_cols] = fixed_df[numeric_cols].round(0)
                # st.write(fixed_df)
                cols = fixed_df.columns.tolist()
                top_level = []
                second_level = []
                for col in cols:
                    c = col.lower()
                    
                    # File Info
                    if 'filename' in c or 'projectname' in c or 'location' in c:
                        top_level.append('File Info')
                        second_level.append(col.replace('-', ' '))
                    
                    elif 'floor' in c:
                        top_level.append('Floor Area (ft²)')
                        second_level.append(col.replace('Floor', ' ').replace('(sqft)', ''))

                    elif 'wall' in c and 'window' not in c:
                        top_level.append('Wall Area (ft²)')
                        second_level.append(col.replace('(sqft)', ' ').replace('(SQFT)', '').replace("Area",'').replace('-',' ').replace("North Wall", "N").replace("East Wall", "E").replace("South Wall", "S").replace("West Wall", "W").replace("SE Wall", "SE").replace("NE Wall", "NE").replace("SW Wall", "SW").replace("NW Wall", "NW").replace("Wall","").replace("Roof","").replace("s+","Above Grade (Wall+Roof)").replace("Total","").replace(" s","t").replace("Allt", "AboveGrade").replace("Above Grade", "Above Grade(Wall+Wind)").replace("AboveGrade", "Above Grade"))
                    
                    elif 'window' in c or 'walls window' in c or 'walls+roof' in c:
                        top_level.append('Window Area (ft²)')
                        second_level.append(col.replace('Window', ' ').replace('(sqft)', '').replace('(SQFT)', '').replace('-','').replace("Area",'').replace("North", "N").replace("South", "S").replace("East", "E").replace("West", "W").replace("  ", " ").replace("Total",""))

                    else:
                        top_level.append('Roof Area (ft²)')
                        second_level.append(col.replace(' ', ' ').replace('(sqft)', '').replace('Area', '').replace("Roof",""))

                # Create MultiIndex
                multi_columns = pd.MultiIndex.from_tuples(list(zip(top_level, second_level)))
                fixed_df.columns = multi_columns
                pd.set_option('display.max_columns', None)
                st.write(fixed_df)
        with colB:
            with st.expander("Detailed Building Information"):
                cData["FileName"] = cData["FileName"].str.replace(r'_[^_]+(?=\.[^.]+$|$)', f'_{user_nm}', regex=True)
                cData = cData.drop(cData.columns[1:8], axis=1)
                # Get original columns
                variable_df = variable_df.drop(variable_df.columns[0], axis=1)
                variable_df = variable_df.drop(['Light(W/Sqft)', 'Equip(W/Sqft)'], axis=1)
                round_dict = {
                    'Power Lighting Total(W)': 0,  # round to 1 decimal
                    'Equipment-Total(W)': 0,  # round to 2 decimals
                    'EFLH': 0,
                    'TOTAL-LOAD(KW)': 0,
                    'LIGHT TO SPACE': 0,
                    'EQUIPMENT TO SPACE': 0,
                    'OCCUPANTS TO SPACE': 0,
                    'WALL CONDUCTION': 0,
                    'ROOF CONDICTION': 0,
                    'WINDOW GLASS+FRM COND': 0,
                    'WINDOW GLASS SOLAR': 0,
                    'ROOF CONDUCTION': 0
                }
                variable_df = variable_df.round(round_dict)

                # st.write(variable_df)
                cols = variable_df.columns.tolist()
                top_level = []
                second_level = []
                for col in cols:
                    c = col.lower()
                    
                    # File Info
                    if 'filename' in c or 'projectname' in c or 'location' in c:
                        top_level.append('File Info')
                        second_level.append(col.replace('-', ' '))
                    
                    # Floor area
                    elif 'floor' in c:
                        top_level.append('Floor Area (ft²)')
                        second_level.append(col.split('(')[0].replace('Floor-Total-', '').replace('-', ' '))
                    
                    # Wall area
                    elif 'wall-area' in c and 'window' not in c:
                        top_level.append('Wall Area (ft²)')
                        second_level.append(col.split('-Area')[0].replace('-', ' '))
                    
                    # Roof area
                    elif 'roof-area' in c or 'walls+roofs-area' in c:
                        top_level.append('Roof Area (ft²)')
                        second_level.append(col.split('-AREA')[0].replace('-', ' '))
                    
                    # Window area
                    elif 'window-area' in c:
                        top_level.append('Window Area (ft²)')
                        second_level.append(col.split('-Area')[0].replace('-', ' '))
                    
                    # U-values
                    elif 'u-value' in c:
                        if 'wall' in c and 'window' not in c:
                            top_level.append('Wall U-Value (BTU/hr·ft²·°F)')
                        elif 'roof' in c:
                            top_level.append('Roof U-Value (BTU/hr·ft²·°F)')
                        elif 'window' in c or 'all walls-window' not in c:
                            top_level.append('Window U-Value (BTU/hr·ft²·°F)')
                        else:
                            top_level.append('Other U-Value')
                        second_level.append(col.split('-U-Value')[0].replace('-', ' ').replace('ROOF', 'Roof').replace("RoofS", "Roof").replace("N Wall", "N").replace("S Wall", "S").replace("E Wall", "E").replace("W Wall", "W").replace("W Window", "W").replace("N Window", "N").replace("S Window", "S").replace("E Window", "E").replace("Window", "").replace("Wall","").replace("ALL WALLS","Above Grade").replace("Roof","").replace("WALLS+","Above Grade(Wall+Roof)"))
                    
                    # Lighting & equipment loads
                    elif 'power' in c or 'equipment-total' in c or 'load' in c or 'equip(w/sqft)' in c or 'light(w/sqft)' in c:
                        top_level.append('Power Density (W)')
                        second_level.append(col.replace('(W)', ' ').replace('Total', '').replace('-', '').replace('Power', '').replace('(KW)','').replace('TOTALLOAD','Total-Load').replace('-',' '))
                    
                    # Conduction & solar
                    elif 'conduction' in c or 'solar' in c or 'cond' in c:
                        top_level.append('Conduction & Solar (kW)')
                        second_level.append(col.replace('-', ' '))
                    
                    # Internal gains
                    elif 'occupants' in c or 'light to space' in c or 'equipment to space' in c or 'process to space' in c or 'infiltration' in c:
                        top_level.append('Internal Gains (kW)')
                        second_level.append(col.replace('-', ' '))
                    
                    elif 'sc' in c:
                        top_level.append('SC')
                        second_level.append(col.replace('SC', ' '))
                    
                    elif 'eflh' in c:
                        top_level.append('EFLH(H)')
                        second_level.append(col.replace('EFLH', ''))
                    
                    elif 'wwr' in c:
                        top_level.append('WWR')
                        second_level.append(col.replace('WWR', ' '))

                    elif 'energy' in c:
                        top_level.append('Energy Use(kWh)')
                        second_level.append(col.replace('(KWH)', ' ').replace('Energy_Outcome', ' '))
                    
                    # Other metrics
                    else:
                        top_level.append('R Value')
                        second_level.append(col.replace('R-Value-', ' '))

                # Create MultiIndex
                multi_columns = pd.MultiIndex.from_tuples(list(zip(top_level, second_level)))

                # Assign to dataframe
                variable_df.columns = multi_columns
                # Optional: show all columns
                pd.set_option('display.max_columns', None)
                st.write(variable_df)
        
        # with colC:
        #     with st.expander("Log File"):
        #         # logFile["File Name"] = logFile["File Name"].str.replace(r'_[^_]+\.sim$', f'_{user_nm}.sim', regex=True)
        #         st.write("logFile")

if st.session_state.script_choice == "tool2":
    # --- Read Excel Sheets ---
    exportxlsx = resource_path(os.path.join("database", "AllData.xlsx"))
    wallDB = pd.read_excel(exportxlsx, sheet_name="Wall")
    roofDB = pd.read_excel(exportxlsx, sheet_name="Roof")
    windDB = pd.read_excel(exportxlsx, sheet_name="Glazing")
    wwrDB = pd.read_excel(exportxlsx, sheet_name="WWR")
    lightDB = pd.read_excel(exportxlsx, sheet_name="Light")
    equipDB = pd.read_excel(exportxlsx, sheet_name="Equip")

    col1, col2, col3, col4, col5, col6 = st.columns(6)
    # --- Wall ---
    with col1:
        wall_options = ["As Designed"] + [f"As Designed + R{r}" for r in [2.5, 5, 7.5, 10, 12.5, 15, 17.5, 20, 22.5, 25, 27.5, 30]]
        wall_choice = st.selectbox("Wall", wall_options)

    # --- Roof ---
    with col2:
        roof_options = ["As Designed"] + [f"As Designed + R{r}" for r in [2.5, 5, 7.5, 10, 12.5, 15, 17.5, 20, 22.5, 25, 27.5, 30]]
        roof_choice = st.selectbox("Roof", roof_options)

    # --- Glazing ---
    with col3:
        glazing_values = windDB["GLASS-TYPE"].iloc[1:].dropna().unique()
        cleaned_values = [v.replace("ML_", "").replace("_Others", "").replace("_", ", ") for v in glazing_values]
        glazing_options = ["As Designed"] + cleaned_values
        glazing_choice = st.selectbox("Glazing", glazing_options)

    # --- WWR ---
    with col4:
        wwr_values = wwrDB["WWR"].iloc[0:].dropna().unique()
        wwr_values_pct = ["As Designed"] + [f"{v*10}%" for v in wwr_values if v != 0]
        wwr_choice = st.selectbox("WWR", wwr_values_pct)

    # --- Light ---
    with col5:
        light_values = lightDB["Percent"].iloc[0:].dropna().unique()
        light_values_pct = ["As Designed"] + [f"{v*100:.0f}%" for v in light_values]
        light_choice = st.selectbox("Lighting Saving", light_values_pct)

    # --- Equip ---
    with col6:
        equip_values = equipDB["Percent"].iloc[0:].dropna().unique()
        equip_values_pct = ["As Designed"] + [f"{v*100:.0f}%" for v in equip_values]
        equip_choice = st.selectbox("Equipment Saving", equip_values_pct)

    # --- Note ---
    st.markdown(
        """**Note:** *Options in Glazings are coded as: <b>GlazingType, U-Value, Shading Coefficient, Lighting-Transmittence</b>.*""",
        unsafe_allow_html=True
    )

    # --- Convert selections to numeric codes ---
    def get_index(choice, options):
        """Return numeric code (0 for 'As Designed', else 1, 2, ...)"""
        return options.index(choice)

    wall_code = get_index(wall_choice, wall_options)
    roof_code = get_index(roof_choice, roof_options)
    glazing_code = get_index(glazing_choice, glazing_options)
    wwr_code = get_index(wwr_choice, wwr_values_pct)
    light_code = get_index(light_choice, light_values_pct)
    equip_code = get_index(equip_choice, equip_values_pct)

    # --- Create single-row dataframe ---
    output_df_New = pd.DataFrame([{
        "Batch_ID": 1,
        "Run_ID": 1,
        "Wall": wall_code,
        "Roof": roof_code,
        "Glazing": glazing_code,
        "WWR": wwr_code,
        "Orient": 0,
        "Light": light_code,
        "Equip": equip_code
    }])

    asdes_df = pd.DataFrame([{
        "Batch_ID": 1,
        "Run_ID": 1,
        "Wall": 0,
        "Roof": 0,
        "Glazing": 0,
        "WWR": 0,
        "Orient": 0,
        "Light": 0,
        "Equip": 0
    }])
    output_df_New = pd.concat([output_df_New, asdes_df], ignore_index=True)
    # st.write(output_df_New)

    from plotly.subplots import make_subplots  
    import plotly.graph_objects as go

    if st.button("Simulate 🚀") and uploaded_file is not None and (output_df_New.iloc[0] != output_df_New.iloc[1]).any():
        with st.spinner("⚡ Processing... This may take a few minutes."):
            os.makedirs(output_inp_folder, exist_ok=True)
            new_batch_id = f"{int(time.time())}"  # unique ID

            selected_rows = output_df_New[output_df_New['Batch_ID'] == run_cnt]
            batch_output_folder = os.path.join(output_inp_folder, f"{user_nm}")
            os.makedirs(batch_output_folder, exist_ok=True)
            # st.write(batch_output_folder)
            num = 1
            modified_files = []
            for _, row in selected_rows.iterrows():
                selected_inp = uploaded_file.name
                new_inp_name = f"{row['Wall']}_{row['Roof']}_{row['Glazing']}_{row['Orient']}_{row['Light']}_{row['WWR']}_{row['Equip']}_{selected_inp}"
                # st.write(new_inp_name)
                new_inp_path = os.path.join(batch_output_folder, new_inp_name)
                inp_file_path = os.path.join(inp_folder, selected_inp)
                if not os.path.exists(inp_file_path):
                    st.error(f"File {inp_file_path} not found. Skipping.")
                    continue
                num += 1

                inp_content = wwr.process_window_insertion_workflow(inp_file_path, row["WWR"])
                inp_content = lighting.updateLPD(inp_content, row['Light'])
                inp_content = equip.updateEquipment(inp_content, row['Equip'])
                inp_content = windows.insert_glass_types_multiple_outputs(inp_content, row['Glazing'])
                inp_content = remove_utility(inp_content)
                inp_content = remove_betweenLightEquip(inp_content)
                count = ModifyWallRoof.count_exterior_walls(inp_content)
                if count > 1:
                    inp_content = ModifyWallRoof.fix_walls(inp_content, row["Wall"])
                    inp_content = ModifyWallRoof.fix_roofs(inp_content, row["Roof"])
                    inp_content = insertRoof.removeDuplicates(inp_content)
                    with open(new_inp_path, 'w') as file:
                        file.writelines(inp_content)
                    modified_files.append(new_inp_name)
                else:
                    st.write("No Exterior-Wall Exists!")

            simulate_files = []
            # Copy and run batch script
            if uploaded_file is None:
                st.error("Please upload an INP file before starting the simulation.")
            else:
                try:
                    script_dir = os.path.dirname(os.path.abspath(__file__))
                    shutil.copy(os.path.join(script_dir, "script.bat"), batch_output_folder)
                    inp_files = [f for f in os.listdir(batch_output_folder) if f.lower().endswith(".inp")]
                    for inp_file in inp_files:
                        file_path = os.path.join(batch_output_folder, os.path.splitext(inp_file)[0])
                        subprocess.call(
                            [os.path.join(batch_output_folder, "script.bat"), file_path, weather_path],
                            shell=True
                        )
                        simulate_files.append(inp_file)
                
                    subprocess.call([os.path.join(batch_output_folder, "script.bat"), batch_output_folder, weather_path], shell=True)
                    required_sections = ['BEPS', 'BEPU', 'LS-C', 'LV-B', 'LV-D', 'PS-E', 'SV-A']
                    log_file_path = check_missing_sections(batch_output_folder, required_sections, new_batch_id, user_nm)
                    get_failed_simulation_data(batch_output_folder, log_file_path)
                    clean_folder(batch_output_folder)
                    combinedData = get_files_for_data_extraction(batch_output_folder, log_file_path, new_batch_id, location_id, user_nm, user_input)
                    merged_df = pd.concat([combinedData], ignore_index=True)
                    # st.write(combinedData)
                    # st.write("ok")
                    generated_names = [os.path.splitext(f)[0] for f in modified_files]
                    merged_df = merged_df[merged_df['FileName'].isin(generated_names)]
                    # st.write(merged_df)
                    merged_df.index = ["As Designed", "ECM"]
                    # Columns contributing to energy
                    energy_cols = [
                        "WALL CONDUCTION", "ROOF CONDUCTION", "WINDOW GLASS+FRM COND", "WINDOW GLASS SOLAR",
                        "DOOR CONDUCTION", "INTERNAL SURFACE COND", "UNDERGROUND SURF COND",
                        "OCCUPANTS TO SPACE", "LIGHT TO SPACE", "EQUIPMENT TO SPACE", "PROCESS TO SPACE", "INFILTRATION"
                    ]

                    # --- Extract and filter out zeros ---
                    as_designed = merged_df.loc["As Designed", energy_cols]
                    ecm = merged_df.loc["ECM", energy_cols]
                    as_designed = as_designed[as_designed > 0]
                    ecm = ecm[ecm > 0]
                    m_df = pd.concat([as_designed, ecm], axis=1)
                    m_df = m_df.rename(columns={'As Designed': 'As Designed (kW)', 'ECM': 'ECM (kW)'})
                    m_df.insert(0, 'Parameters', m_df.index)
                    m_df = m_df.loc[:, m_df.columns.notna()]          # drop NaN column names
                    m_df = m_df.loc[:, m_df.columns != ''] 
                    m_df['As Designed (kW)'] = pd.to_numeric(m_df['As Designed (kW)'], errors='coerce')
                    m_df['ECM (kW)'] = pd.to_numeric(m_df['ECM (kW)'], errors='coerce')
                    m_df['% Saving'] = ((m_df['As Designed (kW)'] - m_df['ECM (kW)']) / m_df['As Designed (kW)']) * 100
                    m_df['% Saving'] = m_df['% Saving'].round(1)
                    # st.write(m_df)

                    # --- Combine unique labels ---
                    all_labels = sorted(set(as_designed.index).union(set(ecm.index)))
                    color_map = {
                        "WALL CONDUCTION": "#1f77b4",     # Blue
                        "ROOF CONDUCTION": "#ff0000",     # Red
                        "WINDOW GLASS+FRM COND": "#87ceeb",  # Light Blue
                        "WINDOW GLASS SOLAR": "#ff6347",  # Light Red / Tomato
                        "DOOR CONDUCTION": "#2ecc71",     # Green
                        "INTERNAL SURFACE COND": "#90ee90",  # Light Green
                        "UNDERGROUND SURF COND": "#4169e1",  # Royal Blue
                        "OCCUPANTS TO SPACE": "#ff7f7f",  # Soft Red
                        "LIGHT TO SPACE": "#3cb371",      # Medium Sea Green
                        "EQUIPMENT TO SPACE": "#add8e6",  # Light Blue
                        "PROCESS TO SPACE": "#66cdaa",    # Medium Aquamarine
                        "INFILTRATION": "#ff4500"         # Orange-Red
                    }

                    # --- Helper to get colors in correct order ---
                    def get_colors(labels):
                        return [color_map.get(l, "#cccccc") for l in labels]

                    # --- Create subplots ---
                    fig = make_subplots(
                        rows=1, cols=2, specs=[[{'type': 'domain'}, {'type': 'domain'}]],
                        subplot_titles=('As Designed', 'ECM')
                    )

                    # --- Add first pie (As Designed) ---
                    fig.add_trace(go.Pie(
                        labels=as_designed.index,
                        values=as_designed.values,
                        # name="As Designed",
                        hole=0.45,  # donut hole
                        textinfo='label+percent',  # show label and %
                        textfont=dict(size=12, color="black"),
                        insidetextorientation='radial',
                        hovertemplate='<b>%{label}</b><br>Value: %{value}<br>Percent: %{percent}',
                        marker=dict(colors=get_colors(as_designed.index), line=dict(color='white', width=2))
                    ), 1, 1)

                    # --- Add second pie (ECM) ---
                    fig.add_trace(go.Pie(
                        labels=ecm.index,
                        values=ecm.values,
                        # name="ECM",
                        hole=0.45,  # donut hole
                        textinfo='label+percent',
                        textfont=dict(size=12, color="black"),
                        insidetextorientation='radial',
                        hovertemplate='<b>%{label}</b><br>Value: %{value}<br>Percent: %{percent}',
                        marker=dict(colors=get_colors(ecm.index), line=dict(color='white', width=2))
                    ), 1, 2)
                    st.markdown(
                        "<h5 style='text-align:left; color:red; font-weight:600;'>Gains Summary</h5>",
                        unsafe_allow_html=True
                    )
                    # --- Layout styling ---
                    fig.update_layout(
                        title_text="",
                        title_font=dict(size=20, color="black", family="Arial"),
                        showlegend=True,
                        legend_orientation="h",
                        legend_x=0.5,
                        legend_y=-0.1,
                        legend_xanchor="center",
                        legend_yanchor="top",
                        # legend_title_text="Energy Components",
                        height=600,
                        margin=dict(t=80, b=120),
                        annotations=[
                            dict(text='As Designed', x=0.18, y=0.5, font_size=14, showarrow=False),
                            dict(text='ECM', x=0.82, y=0.5, font_size=14, showarrow=False)
                        ]
                    )

                    st.plotly_chart(fig, use_container_width=True)
                    st.markdown(
                        "<h5 style='text-align:left; color:red; font-weight:600;'>Building Peak Load Components</h5>",
                        unsafe_allow_html=True
                    )
                    st.write(m_df)

                    dfsi = []
                    for file in os.listdir(batch_output_folder):
                        if file.lower().endswith("_bepu.csv"):
                            file_path = os.path.join(batch_output_folder, file)
                            df = pd.read_csv(file_path)
                            dfsi.append(df)
                    # st.write(f"Loaded {len(dfsi)} BEPU CSV files.")

                    # ⚠️ Here's the fix:
                    if dfsi:  # check if list is not empty
                        dfs = dfsi[0]  # first BEPU CSV
                        # To get the ith DataFrame, make sure 'i' exists
                        i = 1  # or any valid index less than len(dfsi)
                        df = dfsi[i]
                    else:
                        st.warning("No _bepu.csv files found.")

                    df = df.drop(index=[0, 1])
                    dfs = dfs.drop(index=[0, 1])
                    # 2️⃣ Convert all numeric columns
                    numeric_cols = df.columns[4:]
                    numeric_cols1 = dfs.columns[4:]
                    df[numeric_cols] = df[numeric_cols].apply(pd.to_numeric, errors='coerce')
                    dfs[numeric_cols1] = dfs[numeric_cols1].apply(pd.to_numeric, errors='coerce')

                    # 3️⃣ Group by 'Units' and sum numeric columns
                    unit_sum_df = df.groupby("BEPU-UNIT")[numeric_cols].sum().reset_index()
                    unit_sum_dfs = dfs.groupby("BEPU-UNIT")[numeric_cols1].sum().reset_index()

                    # st.write(unit_sum_df)
                    # st.write(unit_sum_dfs)

                    # Filter only KWH rows
                    df_asdes_kwh = unit_sum_dfs[unit_sum_dfs["BEPU-UNIT"] == "KWH"]
                    df_ecm_kwh = unit_sum_df[unit_sum_df["BEPU-UNIT"] == "KWH"]

                    # Get all parameter columns (exclude 'Unnamed: 0', 'BEPU-UNIT', 'TOTAL-BEPU')
                    param_cols = [col for col in unit_sum_dfs.columns if col not in ['Unnamed: 0', 'BEPU-UNIT', 'TOTAL-BEPU']]

                    # Create a combined DataFrame
                    combined_df = pd.DataFrame({
                        "Parameters": param_cols,
                        "As Designed (kWh)": df_asdes_kwh[param_cols].iloc[0].values,
                        "ECM (kWh)": df_ecm_kwh[param_cols].iloc[0].values
                    })

                    # Calculate % Saving
                    combined_df["% Saving"] = ((combined_df["As Designed (kWh)"] - combined_df["ECM (kWh)"]) /
                                                combined_df["As Designed (kWh)"] * 100).round(2)

                    # Drop rows where both are 0
                    combined_df = combined_df[(combined_df["As Designed (kWh)"] > 0) | (combined_df["ECM (kWh)"] > 0)]
                    
                    st.markdown(
                        "<h5 style='text-align:left; color:red; font-weight:600;'>Building Utility Performance</h5>",
                        unsafe_allow_html=True
                    )
                    
                    from plotly.subplots import make_subplots
                    import plotly.graph_objects as go

                    # Remove zero columns
                    unit_sum_df = unit_sum_df.loc[:, (unit_sum_df != 0).any(axis=0)]
                    unit_sum_dfs = unit_sum_dfs.loc[:, (unit_sum_dfs != 0).any(axis=0)]

                    # Prepare data for both pies
                    component_columns_1 = [col for col in unit_sum_df.columns if col != "TOTAL-BEPU"]
                    component_columns_2 = [col for col in unit_sum_dfs.columns if col != "TOTAL-BEPU"]

                    contribution1 = unit_sum_df[component_columns_1].sum().reset_index()
                    contribution1.columns = ["Component", "Value"]

                    contribution2 = unit_sum_dfs[component_columns_2].sum().reset_index()
                    contribution2.columns = ["Component", "Value"]

                    # --- Create subplot with shared legend ---
                    fig = make_subplots(rows=1, cols=2, specs=[[{'type':'domain'}, {'type':'domain'}]],
                                        subplot_titles=('As Designed', 'ECM'),
                                        horizontal_spacing=0.05)  # reduce gap between pies

                    # Add both pies
                    fig.add_trace(
                        go.Pie(labels=contribution1["Component"], values=contribution1["Value"], name="As Designed"),
                        1, 1
                    )
                    fig.add_trace(
                        go.Pie(labels=contribution2["Component"], values=contribution2["Value"], name="ECM"),
                        1, 2
                    )

                    # Update traces (larger pies)
                    fig.update_traces(
                        hole=0.35,  # smaller hole = larger visible pie
                        textinfo="percent+label",
                        hoverinfo="label+value+percent",
                        textfont=dict(size=12),
                    )

                    # --- Layout for larger pies and shared legend ---
                    fig.update_layout(
                        showlegend=True,
                        legend=dict(
                            orientation="h",
                            yanchor="bottom",
                            y=-0.15,
                            xanchor="center",
                            x=0.5
                        ),
                        margin=dict(t=50, b=60, l=30, r=30),
                        height=500,  # increase chart height for bigger pies
                    )

                    # --- Display in Streamlit ---
                    st.plotly_chart(fig, use_container_width=True)

                    # st.markdown(
                    #     "<h5 style='text-align:left; color:red; font-weight:600;'>Energy Use Summary by End Use</h5>",
                    #     unsafe_allow_html=True
                    # )
                    # st.write(combined_df)

                    st.markdown(
                        "<h5 style='text-align:left; color:red; font-weight:600;'>Energy Use Comparison</h5>",
                        unsafe_allow_html=True
                    )
                    custom_colors = ["#FF4B4B", "#1E90FF"]
                    col1, col2 = st.columns(2)
                    with col1:
                        # Rename columns temporarily for clean legend labels
                        temp_df = combined_df.rename(
                            columns={
                                "As Designed (kWh)": "As Designed",
                                "ECM (kWh)": "ECM"
                            }
                        )

                        fig1 = px.bar(
                            temp_df,
                            x="Parameters",
                            y=["As Designed", "ECM"],
                            barmode="group",
                            labels={"value": "Energy (kWh)", "Parameters": ""},  # Remove x-axis label
                            color_discrete_sequence=custom_colors
                        )

                        fig1.update_layout(
                            legend=dict(
                                title_text="",           # no legend title
                                orientation="h",          # horizontal legend
                                yanchor="bottom",
                                y=-0.4,                   # below chart
                                xanchor="center",
                                x=0.5,
                                font=dict(size=12)
                            ),
                            xaxis_tickangle=-45,
                            margin=dict(t=20, b=80)
                        )

                        st.plotly_chart(fig1, use_container_width=True)

                    # --- Chart 2: Stacked Bar Chart ---
                    with col2:
                        combined_df = combined_df.drop('% Saving', axis=1)
                        # --- Build Horizontal Stacked Bar Chart manually ---
                        fig = go.Figure()
                        for param in combined_df["Parameters"]:
                            values = combined_df.loc[combined_df["Parameters"] == param, ["As Designed (kWh)", "ECM (kWh)"]].values[0]
                            fig.add_trace(go.Bar(
                                name=param,
                                y=["As Designed", "ECM"],   # remove (kWh)
                                x=values,
                                text=values,
                                textposition='inside',
                                orientation='h'
                            ))

                        fig.update_layout(
                            barmode='stack',
                            legend=dict(
                                orientation="h",
                                yanchor="top",
                                y=-0.25,
                                xanchor="center",
                                x=0.5
                            ),
                            xaxis_title="Energy (kWh)",
                            yaxis_title=""
                        )

                        # --- Show in Streamlit ---
                        st.plotly_chart(fig, use_container_width=True)

                    st.markdown(
                        "<h5 style='text-align:left; color:red; font-weight:600;'>Energy Use Summary by End Use</h5>",
                        unsafe_allow_html=True
                    )
                    # Calculate % Saving
                    combined_df["% Saving"] = ((combined_df["As Designed (kWh)"] - combined_df["ECM (kWh)"]) /
                                                combined_df["As Designed (kWh)"] * 100).round(2)
                    st.write(combined_df)

                    # # 🟢 INSERT ZIP CREATION + DOWNLOAD CODE HERE ↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓↓
                    # import zipfile
                    # import io

                    # # If combined_Data is a DataFrame, save it to CSV first
                    # combined_path = os.path.join(batch_output_folder, "combined_Data.csv")
                    # if isinstance(combined_Data, pd.DataFrame):
                    #     combined_Data.to_csv(combined_path, index=False)

                    # # Create ZIP buffer in memory
                    # zip_buffer = io.BytesIO()
                    # with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zipf:
                    #     # Add all INP files
                    #     for f in os.listdir(batch_output_folder):
                    #         if f.lower().endswith(".inp"):
                    #             zipf.write(os.path.join(batch_output_folder, f), arcname=f)

                    #     # Add all SIM files
                    #     for f in os.listdir(batch_output_folder):
                    #         if f.lower().endswith(".sim"):
                    #             zipf.write(os.path.join(batch_output_folder, f), arcname=f)

                    #     # Add combined_Data CSV
                    #     if os.path.exists(combined_path):
                    #         zipf.write(combined_path, arcname="combined_Data.csv")

                    # zip_buffer.seek(0)

                    # st.success("✅ All files processed successfully!")
                    # st.download_button(
                    #     label="📦 Download All Results (ZIP)",
                    #     data=zip_buffer,
                    #     file_name=f"{user_nm}_Results.zip",
                    #     mime="application/zip"
                    # )

                    # # 🟢 END INSERTION POINT
                
                except Exception as e:
                    print(f"Moving to next")

    else:
        st.info("Select any Paramters to Compare!")

st.markdown('<hr style="border:1px solid red">', unsafe_allow_html=True)
st.image("images/image123456.png", width=2000) 
st.markdown(
        """
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/5.15.4/css/all.min.css">
    <style>
        .footer {
            background-color: #f8f9fa;
            padding: 20px 0;
            color: #333;
            display: flex;
            justify-content: space-between;
            align-items: center;
            text-align: center;
        }
        .footer .logo {
            flex: 1;
        }
        .footer .logo img {
            max-width: 150px;
            height: auto;
        }
        .footer .social-media {
            flex: 2;
        }
        .footer .social-media p {
            margin: 0;
            font-size: 16px;
        }
        .footer .icons {
            margin-top: 10px;
        }
        .footer .icons a {
            margin: 0 10px;
            color: #666;
            text-decoration: none;
            transition: color 0.3s ease;
        }
        .footer .icons a:hover {
            color: #0077b5; /* LinkedIn color as default */
        }
        .footer .icons a .fab {
            font-size: 28px;
        }
        .footer .additional-content {
            margin-top: 10px;
        }
        .footer .additional-content h4 {
            margin: 0;
            font-size: 18px;
            color: #007bff;
        }
        .footer .additional-content p {
            margin: 5px 0;
            font-size: 16px;
        }
    </style>
    <div style="text-align:center; font-size:14px;">
        Email: <a href="mailto:info@edsglobal.com">info@edsglobal.com</a>&nbsp;&nbsp;&nbsp;|&nbsp;&nbsp;&nbsp;
        Phone: +91 . 11 . 4056 8633&nbsp;&nbsp;&nbsp;|&nbsp;&nbsp;&nbsp;
        <a href="https://twitter.com/edsglobal?lang=en" target="_blank"><i class="fab fa-twitter" style="color:#1DA1F2; margin:0 6px;"></i></a>
        <a href="https://www.facebook.com/Environmental.Design.Solutions/" target="_blank"><i class="fab fa-facebook" style="color:#4267B2; margin:0 6px;"></i></a>
        <a href="https://www.instagram.com/eds_global/?hl=en" target="_blank"><i class="fab fa-instagram" style="color:#E1306C; margin:0 6px;"></i></a>
        <a href="https://www.linkedin.com/company/environmental-design-solutions/" target="_blank"><i class="fab fa-linkedin" style="color:#0077b5; margin:0 6px;"></i></a>
    </div>
    """,
    unsafe_allow_html=True
)